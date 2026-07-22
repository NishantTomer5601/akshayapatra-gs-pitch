"""Append proposed slide-3 v3: solid center pie inside a 2x2 quadrant grid —
gaps in the top two panels, solutions in the bottom two, content bordering the pie."""
from deck_lib import *
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

PATH = "AkshayaPatra_GS_Pitch.pptx"
prs = Presentation(PATH)

s = prs.slides.add_slide(prs.slide_layouts[6])
chrome(s, len(prs.slides._sldIdLst),
       "The two gaps — and the $250,000 that closes them",
       "The gaps highlighted in the supply chain, the grant that closes them, and the "
       "solution we build for each — all in one view.",
       takeaway="Two fixable gaps, two targeted builds — $250k of capital becomes "
                "$310k of savings, every year.")
_, tf = tb(s, SW - ML - 3.6, 0.06, 3.6, 0.20)
para(tf, "PROPOSED SLIDE 3 — v3 (pie quadrant)", size=7.5, color=GRAY, italic=True,
     align=PP_ALIGN.RIGHT, first=True, space_after=0)

# ---------------- 2x2 quadrant panels ----------------
PANEL = RGBColor(0xF4, 0xF4, 0xF4)
QX = [ML, 5.08]; QW_ = 4.52
QY = [1.46, 4.02]; QH = [2.46, 2.44]
for gx in range(2):
    for gy in range(2):
        rect(s, QX[gx], QY[gy], QW_, QH[gy], fill=PANEL, line=None)

def qhead(x, y, text, accent, align_right=False):
    rect(s, x, y + 0.03, 0.06, 0.20, fill=accent, line=None)
    _, tf = tb(s, x + 0.15, y, 3.2, 0.26)
    para(tf, text, size=10.5, color=NAVY, bold=True, first=True, space_after=0)

# ---------------- top-left: gap 1 ----------------
qhead(QX[0] + 0.15, QY[0] + 0.12, "GAP 1 — Procurement", RED)
_, tf = tb(s, QX[0] + 0.18, QY[0] + 0.50, 3.05, 1.85)
para(tf, [("•  Manual, gut-feel indenting", True, TEXT),
          (" — blind to attendance swings, exams & holidays", False, TEXT)],
     size=8, first=True, line_spacing=1.04)
para(tf, [("•  Reactive spot-buying", True, TEXT),
          (" at mandi peak prices; no price intelligence", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("•  No feedback loop", True, TEXT),
          (" from schools on actual meal uptake", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("Cost today:  8–10% meals over-produced · +12% spot premium [dummy]",
           True, RED)], size=8, space_after=0, line_spacing=1.04)

# ---------------- top-right: gap 2 ----------------
qhead(QX[1] + 1.30, QY[0] + 0.12, "GAP 2 — Distribution", RED)
_, tf = tb(s, QX[1] + 1.32, QY[0] + 0.50, 3.05, 1.85)
para(tf, [("•  Static, driver-memory routes", True, TEXT),
          (" — late vans mean children stay hungry all day", False, TEXT)],
     size=8, first=True, line_spacing=1.04)
para(tf, [("•  Untracked vessels & utensils", True, TEXT),
          ("; adulteration & diversion risk en route", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("•  No proof-of-delivery", True, TEXT),
          (" — school shortfall disputes go unresolved", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("Cost today:  6.2% miss the lunch window · $35k/yr shrinkage [dummy]",
           True, RED)], size=8, space_after=0, line_spacing=1.04)

# ---------------- bottom-left: solution 1 ----------------
qhead(QX[0] + 0.15, QY[1] + 0.12, "SOLUTION 1 — “Annapurna AI” · $150k", NAVY)
_, tf = tb(s, QX[0] + 0.18, QY[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("ML demand-forecasting & procurement planning across all 68 kitchens",
           False, TEXT)], size=8, first=True, line_spacing=1.04)
para(tf, [("How: ", True, NAVY),
          ("predicts school-level attendance; auto-generates daily indents & forward "
           "buying schedules", False, TEXT)], size=8, line_spacing=1.04)
para(tf, [("Rollout: ", True, GRAY),
          ("3-kitchen pilot (Mo 1–4) → all 68 kitchens by Mo 12", False, GRAY)],
     size=7.8, line_spacing=1.04)
para(tf, [("→  ~70% less food waste", True, GREEN)], size=8.2, line_spacing=1.04)
para(tf, [("→  est. $230k saved / yr [dummy]", True, GREEN)], size=8.2,
     space_after=0, line_spacing=1.04)

# ---------------- bottom-right: solution 2 ----------------
qhead(QX[1] + 1.30, QY[1] + 0.12, "SOLUTION 2 — “Last-Mile Shield” · $100k", GREEN)
_, tf = tb(s, QX[1] + 1.32, QY[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("AI route optimisation + telematics (GPS, cameras, RFID) + geofenced "
           "smart locks", False, TEXT)], size=8, first=True, line_spacing=1.04)
para(tf, [("How: ", True, NAVY),
          ("drops sequenced to each school's lunch bell; doors unlock only near a "
           "registered school", False, TEXT)], size=8, line_spacing=1.04)
para(tf, [("Bonus: ", True, GRAY),
          ("3 vans freed per hub redeployed to reach new schools", False, GRAY)],
     size=7.8, line_spacing=1.04)
para(tf, [("→  99%+ on-time · −18% fleet km", True, GREEN)], size=8.2,
     line_spacing=1.04)
para(tf, [("→  est. $80k saved / yr [dummy]", True, GREEN)], size=8.2,
     space_after=0, line_spacing=1.04)

# ---------------- center: solid pie on a white badge ----------------
ccx, ccy = 5.0, 3.96
badge_d = 2.70
rect(s, ccx - badge_d/2, ccy - badge_d/2, badge_d, badge_d, fill=WHITE,
     line=LTGRAY, line_w=1.0, shape=MSO_SHAPE.OVAL)

cd = CategoryChartData()
cd.categories = ["Solution 1 — $150k", "Solution 2 — $100k"]
cd.add_series('s', [150, 100])
chart_w = 2.55
gf = s.shapes.add_chart(XL_CHART_TYPE.PIE,
                        Inches(ccx - chart_w/2), Inches(ccy - chart_w/2),
                        Inches(chart_w), Inches(chart_w), cd)
ch = gf.chart
ch.font.size = Pt(9); ch.font.name = SANS; ch.has_title = False
ch.has_legend = False
ser = ch.series[0]
for i, c in enumerate([NAVY, GREEN]):
    pt = ser.points[i]
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
    pt.format.line.color.rgb = WHITE; pt.format.line.width = Pt(2)
plot = ch.plots[0]
plot.has_data_labels = True
dl = plot.data_labels
dl.show_value = True
dl.number_format = '"$"0"k"'; dl.number_format_is_linked = False
dl.font.size = Pt(11); dl.font.bold = True; dl.font.color.rgb = WHITE

# make the chart frame transparent so the white badge circle shows through
cs = ch._chartSpace
spPr = cs.makeelement(qn('c:spPr'), {})
spPr.append(spPr.makeelement(qn('a:noFill'), {}))
ln_el = spPr.makeelement(qn('a:ln'), {})
ln_el.append(ln_el.makeelement(qn('a:noFill'), {}))
spPr.append(ln_el)
ext = cs.find(qn('c:externalData'))
if ext is not None:
    ext.addprevious(spPr)
else:
    cs.append(spPr)

# grant caption pill under the pie
pill_w, pill_h = 2.30, 0.40
pill = rect(s, ccx - pill_w/2, ccy + badge_d/2 - 0.02, pill_w, pill_h, fill=WHITE,
            line=LTGRAY, line_w=0.9, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
_, tf = tb(s, ccx - pill_w/2, ccy + badge_d/2 + 0.005, pill_w, pill_h - 0.05,
           anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE $250k GRANT", True, NAVY)], size=8.5, first=True, space_after=0,
     align=PP_ALIGN.CENTER)
para(tf, "60 / 40 — weighted to the costlier gap", size=6.6, color=GRAY,
     align=PP_ALIGN.CENTER, space_after=0, italic=True)

prs.save(PATH)
print("saved; total slides:", len(prs.slides._sldIdLst))
