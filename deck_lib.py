"""Shared design system for the Akshaya Patra x Goldman Sachs pitch deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---------------- palette (from GS brand panel in reference deck) ----------------
NAVY    = RGBColor(0x1B, 0x3E, 0x6E)   # headlines / emphasis
BLUE    = RGBColor(0x72, 0x97, 0xC5)   # GS brand blue 114,151,197
MIDBLUE = RGBColor(0x3B, 0x6E, 0xA5)
LIGHTBLUE = RGBColor(0xDC, 0xE6, 0xF1)
PALEBLUE  = RGBColor(0xEE, 0xF3, 0xF9)
CREAM   = RGBColor(0xF0, 0xEB, 0xE6)
GOLD    = RGBColor(0xF3, 0xC4, 0x3F)   # functional amber 243,196,63
ORANGE  = RGBColor(0xE0, 0x8A, 0x1E)   # Akshaya Patra accent
RED     = RGBColor(0xC2, 0x23, 0x10)   # functional red 194,35,16
GREEN   = RGBColor(0x39, 0x80, 0x25)   # functional green 57,128,37
PALEGREEN = RGBColor(0xEA, 0xF3, 0xE6)
PALERED   = RGBColor(0xFB, 0xEC, 0xea)
PALEGOLD  = RGBColor(0xFD, 0xF6, 0xE0)
MAGENTA = RGBColor(0xA6, 0x42, 0x8C)   # data-viz 166,66,140
TEAL    = RGBColor(0x09, 0x5F, 0x61)
TEXT    = RGBColor(0x33, 0x33, 0x33)
GRAY    = RGBColor(0x72, 0x73, 0x75)   # secondary text 114,115,117
LTGRAY  = RGBColor(0xBF, 0xBF, 0xBF)   # borders
BANDGRAY= RGBColor(0xD9, 0xD9, 0xD9)   # takeaway banner
FILLGRAY= RGBColor(0xF2, 0xF2, 0xF2)   # placeholder fill
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

VIZ = [NAVY, BLUE, MAGENTA, GOLD, TEAL, GREEN, ORANGE, GRAY]

SERIF = "Georgia"
SANS  = "Arial"

SW, SH = 10.0, 7.5        # slide size, inches (4:3)
ML = 0.40                 # left/right margin

# ---------------- low-level helpers ----------------

def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _set_runs(p, runs, size, color, bold, font, italic=False):
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for item in runs:
        if isinstance(item, str):
            t, b, c = item, bold, color
        elif len(item) == 2:
            t, b = item; c = color
        else:
            t, b, c = item
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = font
        r.font.italic = italic

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    return box, tf

def para(tf, runs, size=9, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         space_after=3, space_before=0, font=SANS, first=False, italic=False,
         line_spacing=1.0):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    _set_runs(p, runs, size, color, bold, font, italic)
    return p

def rect(slide, x, y, w, h, fill=WHITE, line=LTGRAY, line_w=0.75, shape=MSO_SHAPE.RECTANGLE,
         radius=None, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    sp.shadow.inherit = False
    return sp

def line_style(shape, color, width=1.0, dash=None, arrow=False, arrow_w='med'):
    """Style a connector/line: color, weight, dash pattern, arrowhead."""
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    ln = shape.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    if arrow:
        t = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': arrow_w, 'len': arrow_w})
        ln.append(t)
    try: shape.shadow.inherit = False
    except Exception: pass

def dotted_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.25, curve=True):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.CURVE if curve else MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line_style(conn, color, width, dash='sysDash', arrow=True)
    return conn

# ---------------- slide chrome ----------------

def chrome(slide, num, title, subtitle, takeaway=None):
    # top rule
    rl = rect(slide, ML, 0.30, SW - 2*ML, 0.026, fill=NAVY, line=None)
    # title
    _, tf = tb(slide, ML, 0.40, SW - 2*ML, 0.50)
    para(tf, title, size=20, color=NAVY, font=SERIF, first=True, space_after=0)
    _, tf = tb(slide, ML, 0.90, SW - 2*ML, 0.48)
    para(tf, subtitle, size=10.5, color=TEXT, first=True, space_after=0, line_spacing=1.0)
    # takeaway banner
    if takeaway:
        band = rect(slide, ML + 0.15, 6.64, SW - 2*ML - 0.3, 0.40, fill=BANDGRAY, line=None)
        _, tf = tb(slide, ML + 0.30, 6.64, SW - 2*ML - 0.6, 0.40, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, takeaway, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
             font=SERIF, first=True, space_after=0, line_spacing=0.95)
    # footer logo
    _, tf = tb(slide, ML, 7.10, 1.4, 0.36)
    para(tf, "Goldman", size=9, color=RGBColor(0x11, 0x11, 0x11), bold=True, font=SERIF,
         first=True, space_after=0, line_spacing=0.9)
    para(tf, "Sachs", size=9, color=RGBColor(0x11, 0x11, 0x11), bold=True, font=SERIF,
         space_after=0, line_spacing=0.9)
    _, tf = tb(slide, SW - ML - 0.5, 7.14, 0.5, 0.3)
    para(tf, str(num), size=9, color=GRAY, align=PP_ALIGN.RIGHT, first=True)

def section(slide, x, y, w, header, accent=NAVY, size=11.5, rule=True):
    """Open section header: accent tick + bold header + hairline rule. No box."""
    rect(slide, x, y + 0.03, 0.06, 0.20, fill=accent, line=None)
    _, tf = tb(slide, x + 0.16, y, w - 0.16, 0.28)
    para(tf, header, size=size, color=NAVY, bold=True, first=True, space_after=0)
    if rule:
        rect(slide, x, y + 0.30, w, 0.012, fill=LTGRAY, line=None)
    return y + 0.40

def vline(slide, x, y1, y2):
    """Faint vertical column divider."""
    return rect(slide, x, y1, 0.012, y2 - y1, fill=RGBColor(0xDD, 0xDD, 0xDD), line=None)

def panel(slide, x, y, w, h, header=None, accent=NAVY, fill=WHITE, header_size=12.5):
    """White bordered panel with accent bar + header. Returns y where content starts."""
    rect(slide, x, y, w, h, fill=fill, line=LTGRAY, line_w=0.75)
    if header is None:
        return y + 0.12
    rect(slide, x + 0.14, y + 0.13, 0.07, 0.24, fill=accent, line=None)
    _, tf = tb(slide, x + 0.32, y + 0.11, w - 0.45, 0.30)
    para(tf, header, size=header_size, color=NAVY, bold=True, first=True, space_after=0)
    return y + 0.48

def img_placeholder(slide, x, y, w, h, label, sublabel=None):
    sp = rect(slide, x, y, w, h, fill=FILLGRAY, line=GRAY, line_w=1.0)
    ln = sp.line._get_or_add_ln()
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(d)
    _, tf = tb(slide, x + 0.06, y, w - 0.12, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, size=8.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=1)
    if sublabel:
        para(tf, sublabel, size=7, color=GRAY, align=PP_ALIGN.CENTER, space_after=0, italic=True)
    return sp

def stat(slide, x, y, w, value, caption, color=NAVY, vsize=19, csize=8):
    _, tf = tb(slide, x, y, w, 0.75)
    para(tf, value, size=vsize, color=color, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=1, font=SANS)
    para(tf, caption, size=csize, color=GRAY, align=PP_ALIGN.CENTER, space_after=0,
         line_spacing=0.95)

# ---------------- charts ----------------

def _style_chart(chart, font_size=8):
    chart.font.size = Pt(font_size)
    chart.font.name = SANS
    chart.font.color.rgb = TEXT
    chart.has_title = False

def donut(slide, x, y, w, h, categories, values, colors, legend=XL_LEGEND_POSITION.RIGHT,
          show_pct=False, labels=True):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series('s', values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    ch = gf.chart
    _style_chart(ch)
    ch.has_legend = True
    ch.legend.position = legend
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(7.5)
    ser = ch.series[0]
    for i, c in enumerate(colors):
        pt = ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
        pt.format.line.color.rgb = WHITE
        pt.format.line.width = Pt(1)
    if labels:
        plot = ch.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(7.5)
        dl.font.bold = True
        dl.font.color.rgb = WHITE
        if show_pct:
            dl.show_percentage = True
            dl.show_value = False
    return ch

def col_chart(slide, x, y, w, h, categories, series, colors, gap=60, legend=False,
              value_fmt=None, font_size=8):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    ch = gf.chart
    _style_chart(ch, font_size)
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(7.5)
    plot = ch.plots[0]
    plot.gap_width = gap
    for i, s in enumerate(ch.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = colors[i % len(colors)]
        s.format.line.fill.background()
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(7.5)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = NAVY
    if value_fmt:
        plot.data_labels.number_format = value_fmt
        plot.data_labels.number_format_is_linked = False
    va = ch.value_axis
    va.has_major_gridlines = False
    va.visible = False
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(7.5)
    ca.format.line.color.rgb = LTGRAY
    return ch

def line_chart(slide, x, y, w, h, categories, series, colors, legend=False, smooth=True,
               font_size=8):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    ch = gf.chart
    _style_chart(ch, font_size)
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(7.5)
    for i, s in enumerate(ch.series):
        s.smooth = smooth
        s.format.line.color.rgb = colors[i % len(colors)]
        s.format.line.width = Pt(2.25)
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    va.major_gridlines.format.line.width = Pt(0.5)
    va.tick_labels.font.size = Pt(7.5)
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(7.5)
    ca.format.line.color.rgb = LTGRAY
    return ch
