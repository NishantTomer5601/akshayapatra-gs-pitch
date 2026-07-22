"""Standalone mock: proposed 'maximum impact' data visual for Solution 1 savings.
Builds a single-slide pptx copy — the real deck is untouched."""
from deck_lib import *
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
import shutil

SRC = "AkshayaPatra_GS_Pitch.pptx"
MOCK = ("/private/tmp/claude-501/-Users-nishanttomer-Desktop/"
        "ad99222a-484a-437c-9678-2fdf9eec1d5d/scratchpad/chart_mock.pptx")

prs = Presentation(SRC)
s = prs.slides.add_slide(prs.slide_layouts[6])
chrome(s, len(prs.slides._sldIdLst), "Solution 1 — where the savings come from, and what they buy",
       "Every bar is a number from TAPF's own data — wastage logs and seasonal price "
       "spreads — built up transparently into recurring annual impact.",
       takeaway="₹6 Cr saved every year = 9M+ more meals = 40,000+ children fed for a "
                "full year — recurring, audited, conservative.")
_, tf = tb(s, SW - ML - 4.0, 0.06, 4.0, 0.20)
para(tf, "PROPOSED SLIDE 4 GRAPH — v2 (compare with slide 9)", size=7.5, color=GRAY, italic=True,
     align=PP_ALIGN.RIGHT, first=True, space_after=0)

# ---------------- waterfall (left, manual shapes) ----------------
cy = section(s, ML, 1.50, 5.60, "Annual savings build-up (₹ Cr / year)", accent=NAVY)
BASE_Y = 5.30          # x-axis line
SCALE = 0.38           # inches per ₹Cr
BW = 0.80              # bar width
xs = [0.62, 1.70, 2.78, 3.86, 4.94]

bars = [
    # (label lines, source chip, from, to, color, value label)
    ("Veg waste cut\n4.9% → 2%", "TAPF wastage logs", 0.0, 4.1, GREEN, "+4.1"),
    ("Dhals buy-timing\n6–7% spread", "TAPF price data", 4.1, 6.0, TEAL, "+1.9"),
    ("Oils buy-timing\n4–5% spread", "TAPF price data", 6.0, 6.8, GOLD, "+0.8"),
    ("Storage &\ncarry cost", "netted off", 6.8, 6.0, RED, "−0.8"),
    ("NET SAVINGS\nevery year", "recurring", 0.0, 6.0, NAVY, "₹6.0 Cr"),
]
# axis line
rect(s, ML, BASE_Y, 5.55, 0.014, fill=LTGRAY, line=None)
prev_top = None
for (lbl, chip, v0, v1, c, vlab), x in zip(bars, xs):
    lo, hi = min(v0, v1), max(v0, v1)
    y_top = BASE_Y - hi * SCALE
    h = (hi - lo) * SCALE
    rect(s, x, y_top, BW, h, fill=c, line=None)
    # value label above bar
    _, tf = tb(s, x - 0.15, y_top - 0.24, BW + 0.30, 0.22)
    para(tf, vlab, size=9.5, color=(RED if v1 < v0 else NAVY), bold=True,
         align=PP_ALIGN.CENTER, first=True, space_after=0)
    # x labels + source chip
    _, tf = tb(s, x - 0.16, BASE_Y + 0.06, BW + 0.32, 0.55)
    para(tf, lbl, size=6.9, color=TEXT, bold=True, align=PP_ALIGN.CENTER,
         first=True, space_after=1, line_spacing=0.9)
    para(tf, chip, size=6, color=GRAY, italic=True, align=PP_ALIGN.CENTER,
         space_after=0)
    # dotted connector from previous bar top
    if prev_top is not None:
        lvl = BASE_Y - v0 * SCALE if v1 >= v0 else BASE_Y - v0 * SCALE
        conn = s.shapes.add_connector(1, Inches(prev_x + BW), Inches(prev_lvl),
                                      Inches(x), Inches(prev_lvl))
        line_style(conn, GRAY, 0.9, dash='sysDash')
    prev_top = y_top
    prev_x = x
    prev_lvl = BASE_Y - v1 * SCALE

# ---------------- right top: ramp mini-chart ----------------
RX, RW_ = 6.42, SW - ML - 6.42
cy = section(s, RX, 1.50, RW_, "Capture ramp (₹ Cr saved / yr)", accent=MIDBLUE,
             size=9.5)
cd = CategoryChartData()
cd.categories = ["2026", "2027", "2028", "2029"]
cd.add_series("saved", [1.8, 3.6, 5.1, 5.1])
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(RX), Inches(cy),
                        Inches(RW_), Inches(1.30), cd)
ch = gf.chart
ch.font.size = Pt(6.5); ch.font.name = SANS; ch.font.color.rgb = TEXT
ch.has_title = False; ch.has_legend = False
ch.series[0].format.fill.solid()
ch.series[0].format.fill.fore_color.rgb = GREEN
ch.series[0].format.line.fill.background()
plot = ch.plots[0]; plot.gap_width = 60
plot.has_data_labels = True
plot.data_labels.font.size = Pt(6.5); plot.data_labels.font.bold = True
plot.data_labels.font.color.rgb = NAVY
va = ch.value_axis; va.has_major_gridlines = False; va.visible = False
ca = ch.category_axis; ca.tick_labels.font.size = Pt(6.5)
ca.format.line.color.rgb = LTGRAY
_, tf = tb(s, RX, cy + 1.34, RW_, 0.20)
para(tf, "pilot 30% → scale 60% → steady-state 85% capture", size=6.5, color=GRAY,
     italic=True, first=True, space_after=0, align=PP_ALIGN.CENTER)

# ---------------- right bottom: conversion cascade ----------------
cy2 = section(s, RX, 3.62, RW_, "What ₹6 Cr/yr buys — every year", accent=GREEN,
              size=9.5)
steps = [("₹6.0 Cr", "saved annually, steady-state", NAVY),
         ("at ₹6.56", "donor cost per meal (govt funds the rest)", GRAY),
         ("9.1M+", "additional meals every year", GREEN),
         ("40,000+", "children fed for a full year — every year", ORANGE)]
sy = cy2 + 0.02
for i, (v, c, col) in enumerate(steps):
    _, tf = tb(s, RX, sy, RW_, 0.56)
    para(tf, [(v + "  ", True, col), (c, False, TEXT)], size=9.5 if i != 1 else 8,
         first=True, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, RX + 0.30, sy + 0.28, RX + 0.30, sy + 0.52, color=LTGRAY,
                     curve=False, width=1.0)
    sy += 0.60

# ---------------- assumptions footnote ----------------
_, tf = tb(s, ML, 6.10, 5.60, 0.44)
para(tf, "Assumptions: illustrative commodity split of TAPF's ₹450 Cr+ food spend; "
         "seasonal spreads per TAPF (dhals 6–7%, oils 4–5%, excl. taxation & global "
         "events); wastage per TAPF supply-chain data (4.1–5.6%/month). Spatial "
         "sourcing & spot-premium avoidance: not modelled — upside.",
     size=6.2, color=GRAY, italic=True, first=True, space_after=0, line_spacing=1.0)
_, tf = tb(s, 5.30, 7.16, 3.72, 0.30)
para(tf, "Source: The Akshaya Patra Foundation — official documents", size=6.5,
     color=GRAY, italic=True, align=PP_ALIGN.RIGHT, first=True, space_after=0)

prs.save(SRC)
print("saved; total slides:", len(prs.slides._sldIdLst))
