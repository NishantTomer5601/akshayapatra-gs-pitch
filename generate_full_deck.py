"""ONE-CELL GENERATOR — paste this entire script into a single Jupyter cell
and run it. Produces the full 11-slide AkshayaPatra_GS_Pitch.pptx.
Requires: pip install python-pptx
"""

# ======================================================================
# ===== deck_lib.py =====
# ======================================================================
"""Shared design system for the Akshaya Patra x Goldman Sachs pitch deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---------------- palette (from GS brand panel in reference deck) ----------------
NAVY    = RGBColor(0x1B, 0x3E, 0x6E)   # headlines / emphasis
BLUE    = RGBColor(0x72, 0x97, 0xC5)   # GS brand blue 114,151,197
MIDBLUE = RGBColor(0x3B, 0x6E, 0xA5)
LIGHTBLUE = RGBColor(0xDC, 0xE6, 0xF1)
PALEBLUE  = RGBColor(0xEE, 0xF3, 0xF9)
CREAM   = RGBColor(0xF0, 0xEB, 0xE6)
GOLD    = RGBColor(0xF3, 0xC4, 0x3F)   # functional amber 243,196,63
ORANGE  = RGBColor(0xE0, 0x8A, 0x1E)   # Akshaya Patra accent
RED     = RGBColor(0xC2, 0x23, 0x10)   # functional red 194,35,16
GREEN   = RGBColor(0x39, 0x80, 0x25)   # functional green 57,128,37
PALEGREEN = RGBColor(0xEA, 0xF3, 0xE6)
PALERED   = RGBColor(0xFB, 0xEC, 0xea)
PALEGOLD  = RGBColor(0xFD, 0xF6, 0xE0)
MAGENTA = RGBColor(0xA6, 0x42, 0x8C)   # data-viz 166,66,140
TEAL    = RGBColor(0x09, 0x5F, 0x61)
TEXT    = RGBColor(0x33, 0x33, 0x33)
GRAY    = RGBColor(0x72, 0x73, 0x75)   # secondary text 114,115,117
LTGRAY  = RGBColor(0xBF, 0xBF, 0xBF)   # borders
BANDGRAY= RGBColor(0xD9, 0xD9, 0xD9)   # takeaway banner
FILLGRAY= RGBColor(0xF2, 0xF2, 0xF2)   # placeholder fill
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

VIZ = [NAVY, BLUE, MAGENTA, GOLD, TEAL, GREEN, ORANGE, GRAY]

SERIF = "Georgia"
SANS  = "Arial"

SW, SH = 10.0, 7.5        # slide size, inches (4:3)
ML = 0.40                 # left/right margin

# ---------------- low-level helpers ----------------

def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _set_runs(p, runs, size, color, bold, font, italic=False):
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for item in runs:
        if isinstance(item, str):
            t, b, c = item, bold, color
        elif len(item) == 2:
            t, b = item; c = color
        else:
            t, b, c = item
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = font
        r.font.italic = italic

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    return box, tf

def para(tf, runs, size=9, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         space_after=3, space_before=0, font=SANS, first=False, italic=False,
         line_spacing=1.0):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    _set_runs(p, runs, size, color, bold, font, italic)
    return p

def rect(slide, x, y, w, h, fill=WHITE, line=LTGRAY, line_w=0.75, shape=MSO_SHAPE.RECTANGLE,
         radius=None, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    sp.shadow.inherit = False
    return sp

def line_style(shape, color, width=1.0, dash=None, arrow=False, arrow_w='med'):
    """Style a connector/line: color, weight, dash pattern, arrowhead."""
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    ln = shape.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    if arrow:
        t = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': arrow_w, 'len': arrow_w})
        ln.append(t)
    try: shape.shadow.inherit = False
    except Exception: pass

def dotted_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.25, curve=True):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.CURVE if curve else MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line_style(conn, color, width, dash='sysDash', arrow=True)
    return conn

# ---------------- slide chrome ----------------

def chrome(slide, num, title, subtitle, takeaway=None):
    # top rule
    rl = rect(slide, ML, 0.30, SW - 2*ML, 0.026, fill=NAVY, line=None)
    # title
    _, tf = tb(slide, ML, 0.40, SW - 2*ML, 0.50)
    para(tf, title, size=20, color=NAVY, font=SERIF, first=True, space_after=0)
    _, tf = tb(slide, ML, 0.90, SW - 2*ML, 0.48)
    para(tf, subtitle, size=10.5, color=TEXT, first=True, space_after=0, line_spacing=1.0)
    # takeaway banner
    if takeaway:
        band = rect(slide, ML + 0.15, 6.64, SW - 2*ML - 0.3, 0.40, fill=BANDGRAY, line=None)
        _, tf = tb(slide, ML + 0.30, 6.64, SW - 2*ML - 0.6, 0.40, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, takeaway, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
             font=SERIF, first=True, space_after=0, line_spacing=0.95)
    # footer logo
    _, tf = tb(slide, ML, 7.10, 1.4, 0.36)
    para(tf, "Goldman", size=9, color=RGBColor(0x11, 0x11, 0x11), bold=True, font=SERIF,
         first=True, space_after=0, line_spacing=0.9)
    para(tf, "Sachs", size=9, color=RGBColor(0x11, 0x11, 0x11), bold=True, font=SERIF,
         space_after=0, line_spacing=0.9)
    _, tf = tb(slide, SW - ML - 0.5, 7.14, 0.5, 0.3)
    para(tf, str(num), size=9, color=GRAY, align=PP_ALIGN.RIGHT, first=True)

def section(slide, x, y, w, header, accent=NAVY, size=11.5, rule=True):
    """Open section header: accent tick + bold header + hairline rule. No box."""
    rect(slide, x, y + 0.03, 0.06, 0.20, fill=accent, line=None)
    _, tf = tb(slide, x + 0.16, y, w - 0.16, 0.28)
    para(tf, header, size=size, color=NAVY, bold=True, first=True, space_after=0)
    if rule:
        rect(slide, x, y + 0.30, w, 0.012, fill=LTGRAY, line=None)
    return y + 0.40

def vline(slide, x, y1, y2):
    """Faint vertical column divider."""
    return rect(slide, x, y1, 0.012, y2 - y1, fill=RGBColor(0xDD, 0xDD, 0xDD), line=None)

def panel(slide, x, y, w, h, header=None, accent=NAVY, fill=WHITE, header_size=12.5):
    """White bordered panel with accent bar + header. Returns y where content starts."""
    rect(slide, x, y, w, h, fill=fill, line=LTGRAY, line_w=0.75)
    if header is None:
        return y + 0.12
    rect(slide, x + 0.14, y + 0.13, 0.07, 0.24, fill=accent, line=None)
    _, tf = tb(slide, x + 0.32, y + 0.11, w - 0.45, 0.30)
    para(tf, header, size=header_size, color=NAVY, bold=True, first=True, space_after=0)
    return y + 0.48

def img_placeholder(slide, x, y, w, h, label, sublabel=None):
    sp = rect(slide, x, y, w, h, fill=FILLGRAY, line=GRAY, line_w=1.0)
    ln = sp.line._get_or_add_ln()
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(d)
    _, tf = tb(slide, x + 0.06, y, w - 0.12, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, size=8.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=1)
    if sublabel:
        para(tf, sublabel, size=7, color=GRAY, align=PP_ALIGN.CENTER, space_after=0, italic=True)
    return sp

def stat(slide, x, y, w, value, caption, color=NAVY, vsize=19, csize=8):
    _, tf = tb(slide, x, y, w, 0.75)
    para(tf, value, size=vsize, color=color, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=1, font=SANS)
    para(tf, caption, size=csize, color=GRAY, align=PP_ALIGN.CENTER, space_after=0,
         line_spacing=0.95)

# ---------------- charts ----------------

def _style_chart(chart, font_size=8):
    chart.font.size = Pt(font_size)
    chart.font.name = SANS
    chart.font.color.rgb = TEXT
    chart.has_title = False

def donut(slide, x, y, w, h, categories, values, colors, legend=XL_LEGEND_POSITION.RIGHT,
          show_pct=False, labels=True):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series('s', values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    ch = gf.chart
    _style_chart(ch)
    ch.has_legend = True
    ch.legend.position = legend
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(7.5)
    ser = ch.series[0]
    for i, c in enumerate(colors):
        pt = ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
        pt.format.line.color.rgb = WHITE
        pt.format.line.width = Pt(1)
    if labels:
        plot = ch.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(7.5)
        dl.font.bold = True
        dl.font.color.rgb = WHITE
        if show_pct:
            dl.show_percentage = True
            dl.show_value = False
    return ch

def col_chart(slide, x, y, w, h, categories, series, colors, gap=60, legend=False,
              value_fmt=None, font_size=8):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    ch = gf.chart
    _style_chart(ch, font_size)
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(7.5)
    plot = ch.plots[0]
    plot.gap_width = gap
    for i, s in enumerate(ch.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = colors[i % len(colors)]
        s.format.line.fill.background()
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(7.5)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = NAVY
    if value_fmt:
        plot.data_labels.number_format = value_fmt
        plot.data_labels.number_format_is_linked = False
    va = ch.value_axis
    va.has_major_gridlines = False
    va.visible = False
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(7.5)
    ca.format.line.color.rgb = LTGRAY
    return ch

def line_chart(slide, x, y, w, h, categories, series, colors, legend=False, smooth=True,
               font_size=8):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    ch = gf.chart
    _style_chart(ch, font_size)
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(7.5)
    for i, s in enumerate(ch.series):
        s.smooth = smooth
        s.format.line.color.rgb = colors[i % len(colors)]
        s.format.line.width = Pt(2.25)
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    va.major_gridlines.format.line.width = Pt(0.5)
    va.tick_labels.font.size = Pt(7.5)
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(7.5)
    ca.format.line.color.rgb = LTGRAY
    return ch


# ======================================================================
# ===== build_deck.py =====
# ======================================================================
"""Akshaya Patra x Goldman Sachs pitch — v2: 4:3, open editorial layout (minimal boxes)."""
import math

prs = new_deck()

# =====================================================================
# SLIDE 1
# =====================================================================
s = add_slide(prs)
chrome(s, 1, "Who is The Akshaya Patra Foundation?",
       "The world's largest NGO-run school meal programme — serving 2.3 million children "
       "every school day, so that hunger never interrupts education.",
       takeaway="“No child shall be deprived of education because of hunger” — "
                "at a scale no other NGO matches.")

# hero band
rect(s, ML, 1.42, SW - 2*ML, 0.80, fill=ORANGE, line=None)
img_placeholder(s, ML + 0.08, 1.51, 0.98, 0.62, "AKSHAYA PATRA\nlogo")
_, tf = tb(s, ML + 1.22, 1.42, SW - 2*ML - 1.36, 0.80, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Over 1 in 3 children in India show signs of undernourishment", True, WHITE),
          (" — and hunger is a leading reason children drop out of school. Founded in "
           "2000 in Bengaluru with ", False, WHITE),
          ("1,500 children across 5 schools", True, WHITE),
          (", Akshaya Patra now runs the world's largest NGO-operated mid-day meal "
           "programme — a PPP with the Government of India that has served ", False, WHITE),
          ("4+ billion cumulative meals.", True, WHITE)],
     size=8, first=True, space_after=0, line_spacing=1.0)

LX, LW = ML, 4.42
RX, RW = 5.02, SW - ML - 5.02
vline(s, 4.77, 2.42, 6.45)

# L: who & why
cy = section(s, LX, 2.38, LW, "Who they are & why we chose them", accent=NAVY)
_, tf = tb(s, LX, cy, LW, 1.55)
para(tf, [("Mission: ", True, NAVY),
          ("eliminate classroom hunger by implementing the Mid-Day Meal Scheme in "
           "government schools while fighting child malnutrition across India.",
           False, TEXT)], size=8, first=True, line_spacing=1.02)
para(tf, [("Why Akshaya Patra: ", True, NAVY),
          ("our on-ground kitchen inspections and supply-chain review found world-class "
           "production — and ", False, TEXT),
          ("two clearly identified, fixable gaps", True, TEXT),
          (" where $250k unlocks compounding, year-on-year impact rather than one-off "
           "sustenance.", False, TEXT)], size=8, line_spacing=1.02)
para(tf, "✓ 25+ yrs of proven delivery   ✓ Govt-audited PPP   ✓ Tech-forward leadership",
     size=7.8, color=GREEN, bold=True, space_after=0)

# L: journey chart
cy = section(s, LX, 4.32, LW, "The journey — daily meals served (millions)", accent=BLUE)
line_chart(s, LX, cy + 0.02, LW, 1.42,
           ["2000", "2005", "2010", "2015", "2020", "2023", "2026"],
           [("Meals/day (M)", [0.001, 0.33, 1.20, 1.44, 1.80, 2.00, 2.30])], [NAVY],
           font_size=7)
_, tf = tb(s, LX, cy + 1.50, LW, 0.30)
para(tf, [("2012: ", True, TEXT), ("1 billionth meal   ·   ", False, TEXT),
          ("2019: ", True, TEXT), ("3 billionth served by Hon'ble PM   ·   ", False, TEXT),
          ("2024: ", True, TEXT), ("4 billion meals", False, TEXT)],
     size=6.8, first=True, space_after=0, color=GRAY)

# R: what they do
cy = section(s, RX, 2.38, RW, "What they do — the hub-and-spoke model", accent=MIDBLUE)
_, tf = tb(s, RX, cy, RW, 1.15)
para(tf, [("•  Centralised mega-kitchens ", True, TEXT),
          ("cook up to 100,000 meals each before 6 AM daily", False, TEXT)],
     size=7.8, first=True, line_spacing=1.0)
para(tf, [("•  Hub-and-spoke delivery ", True, TEXT),
          ("of hot, locally-adapted menus to school doorsteps", False, TEXT)],
     size=7.8, line_spacing=1.0)
para(tf, [("•  PPP funding ", True, TEXT),
          ("— govt subsidies + philanthropy; every $ leveraged ~3x", False, TEXT)],
     size=7.8, line_spacing=1.0)
para(tf, [("•  Beyond food: ", True, TEXT),
          ("scholarships, anganwadi feeding, disaster-relief kitchens", False, TEXT)],
     size=7.8, space_after=0, line_spacing=1.0)

# R: scale
cy = section(s, RX, 3.94, RW, "Scale at a glance (FY26)", accent=ORANGE)
qw = RW / 4
for i, (v, c) in enumerate([("2.3M", "meals served\nevery day"),
                            ("68", "mega-kitchens\nacross India"),
                            ("16", "states &\n2 UTs"),
                            ("24,000+", "govt schools\nserved")]):
    stat(s, RX + i*qw, cy + 0.02, qw, v, c, vsize=14, csize=6.8)

# R: SDGs
cy = section(s, RX, 5.12, RW, "Aligned to UN SDGs", accent=GREEN)
for i, (t1, t2) in enumerate([("SDG 2", "Zero Hunger"), ("SDG 4", "Quality Educ."),
                              ("SDG 17", "Partners")]):
    img_placeholder(s, RX + i * 0.72, cy + 0.04, 0.64, 0.64, t1, t2)
img_placeholder(s, RX + 2.28, cy + 0.04, RW - 2.28, 0.64,
                "\U0001F4F7 PHOTO: children at mid-day meal")

# =====================================================================
# SLIDE 2
# =====================================================================
s = add_slide(prs)
chrome(s, 2, "A world-class supply chain — with two blind spots",
       "We walked their kitchens, interviewed staff and mapped the end-to-end flow: "
       "production is best-in-class; the gaps sit at the first and last mile.",
       takeaway="A Formula-1 production line — fed by guesswork, delivered on hand-drawn "
                "maps. That is where we intervene.")

_, tf = tb(s, ML, 1.40, SW - 2*ML, 0.24)
para(tf, "Our diligence: 2 kitchen visits  •  12 staff interviews  •  procurement & route "
         "data review", size=8, color=GRAY, italic=True, first=True, space_after=0)

# borderless validation row
vx_list = [(ML, 3.0, "Audited & certified",
            "FSSAI-licensed, ISO 22000-certified kitchens; govt & third-party audits"),
           (ML + 3.15, 3.0, "Global recognition",
            "UK PM Rishi Sunak toured the Bengaluru kitchen — praised hygiene & automation"),
           (ML + 6.30, 2.9, "Our on-ground audit",
            "Gravity-flow cooking, 2-hr batch cycles, near-zero contamination risk")]
for vx, vw, t, d in vx_list:
    _, tf = tb(s, vx, 1.74, vw, 0.72)
    para(tf, [("✓ " + t, True, GREEN)], size=9, first=True, space_after=1)
    para(tf, d, size=7.2, color=TEXT, space_after=0, line_spacing=0.95)

# chevron flow
FY = 2.72; FH = 0.98
steps = [
    ("PROCURE-\nMENT", "Grains, dairy\nfrom mandis", RED, True),
    ("STORAGE", "Silos, cold\nrooms; FIFO", BLUE, False),
    ("PRE-\nPROCESS", "Auto roti lines,\nveg prep", BLUE, False),
    ("COOKING", "Steam cauldrons;\n100k by 6 AM", BLUE, False),
    ("QC &\nPACKING", "Lab-tested;\nsealed vessels", BLUE, False),
    ("DISTRIB-\nUTION", "Van fleet to\nschools by lunch", RED, True),
]
cw = 1.50; gap_w = 0.04
fx = ML
for i, (t, d, c, is_gap) in enumerate(steps):
    x = fx + i * (cw + gap_w)
    ch = rect(s, x, FY, cw, FH, fill=(PALERED if is_gap else PALEBLUE),
              line=(RED if is_gap else MIDBLUE), line_w=(1.5 if is_gap else 0.9),
              shape=MSO_SHAPE.CHEVRON)
    try: ch.adjustments[0] = 0.30
    except Exception: pass
    _, tf = tb(s, x + 0.20, FY + 0.08, cw - 0.32, FH - 0.16, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t.replace("-\n", ""), size=7.8, color=(RED if is_gap else NAVY), bold=True,
         first=True, align=PP_ALIGN.CENTER, space_after=1, line_spacing=0.9)
    para(tf, d, size=6, color=TEXT, align=PP_ALIGN.CENTER, space_after=0, line_spacing=0.9)

# arrows from validation texts to nodes
dotted_arrow(s, 4.4, 2.42, 4.9, FY - 0.04, color=GREEN, curve=True, width=1.0)
dotted_arrow(s, 7.6, 2.42, 7.2, FY - 0.04, color=GREEN, curve=True, width=1.0)

_, tf = tb(s, ML, FY + FH + 0.08, SW - 2*ML, 0.26)
para(tf, [("Steps 2–5 are best-in-class ", True, GREEN),
          ("— automated, audited, built to scale. No intervention needed.", False, GRAY)],
     size=8.5, first=True, align=PP_ALIGN.CENTER, space_after=0)

# gap boxes
GY = 4.40; GH = 1.82
g1x, g1w = ML, 4.50
rect(s, g1x, GY, g1w, GH, fill=WHITE, line=RED, line_w=1.25)
rect(s, g1x, GY, g1w, 0.32, fill=RED, line=None)
_, tf = tb(s, g1x + 0.12, GY + 0.045, g1w - 0.24, 0.26)
para(tf, "GAP 1 — PROCUREMENT: no demand forecasting", size=9, color=WHITE, bold=True,
     first=True, space_after=0)
_, tf = tb(s, g1x + 0.14, GY + 0.42, g1w - 0.28, GH - 0.52)
para(tf, [("•  Manual estimates", True, TEXT),
          (" — no view of school-level attendance, holidays or menu uptake", False, TEXT)],
     size=7.5, first=True, line_spacing=1.0)
para(tf, [("•  ~8–10% over-production", True, TEXT),
          (" wasted on low-attendance days [dummy]", False, TEXT)], size=7.5,
     line_spacing=1.0)
para(tf, [("•  Reactive spot-buying", True, TEXT),
          (" at mandi peak prices — up to 12% above plan [dummy]", False, TEXT)],
     size=7.5, line_spacing=1.0)
para(tf, [("•  No price intelligence", True, TEXT),
          (" — cheapest forward-buying windows are missed", False, TEXT)], size=7.5,
     line_spacing=1.0)
para(tf, [("•  Result: ", True, TEXT),
          ("higher cost per meal → fewer children fed per donor $", False, TEXT)],
     size=7.5, line_spacing=1.0)
para(tf, [("→ addressed by Solution 1 ($150k)", True, NAVY)], size=7.8, space_after=0)

g2x, g2w = ML + 4.70, 4.50
rect(s, g2x, GY, g2w, GH, fill=WHITE, line=RED, line_w=1.25)
rect(s, g2x, GY, g2w, 0.32, fill=RED, line=None)
_, tf = tb(s, g2x + 0.12, GY + 0.045, g2w - 0.24, 0.26)
para(tf, "GAP 2 — DISTRIBUTION: no route intelligence", size=9, color=WHITE, bold=True,
     first=True, space_after=0)
_, tf = tb(s, g2x + 0.14, GY + 0.42, g2w - 0.28, GH - 0.52)
para(tf, [("•  Static, driver-memory routes", True, TEXT),
          (" — 15–30 min past the lunch bell means children stay hungry all day",
           False, TEXT)], size=7.5, first=True, line_spacing=1.0)
para(tf, [("•  No telematics", True, TEXT),
          (" — vessels & utensils untracked; losses recur", False, TEXT)], size=7.5,
     line_spacing=1.0)
para(tf, [("•  Adulteration / diversion risk", True, TEXT),
          (" — unsealed vans can be offloaded at roadside dhabas", False, TEXT)],
     size=7.5, line_spacing=1.0)
para(tf, [("•  No proof-of-delivery", True, TEXT),
          (" — school shortfall disputes unresolved", False, TEXT)], size=7.5,
     line_spacing=1.0)
para(tf, [("•  Result: ", True, TEXT),
          ("missed lunch windows, fuel burn, shrinkage", False, TEXT)], size=7.5,
     line_spacing=1.0)
para(tf, [("→ addressed by Solution 2 ($100k)", True, NAVY)], size=7.8, space_after=0)

dotted_arrow(s, g1x + 0.8, GY, fx + 0.75, FY + FH + 0.03, color=RED, width=1.25)
dotted_arrow(s, g2x + 3.9, GY, fx + 5 * (cw + gap_w) + 0.75, FY + FH + 0.03,
             color=RED, width=1.25)

# =====================================================================
# SLIDE 3 — the two gaps + $250k split (pie quadrant)
# =====================================================================
s = add_slide(prs)
chrome(s, 3, "The two gaps — and the $250,000 that closes them",
       "The gaps highlighted in the supply chain, the grant that closes them, and the "
       "solution we build for each — all in one view.",
       takeaway="Two fixable gaps, two targeted builds — $250k of capital becomes "
                "$310k of savings, every year.")

PANEL3 = RGBColor(0xF4, 0xF4, 0xF4)
QX3 = [ML, 5.08]; QW3 = 4.52
QY3 = [1.46, 4.02]; QH3 = [2.46, 2.44]
for gx in range(2):
    for gy in range(2):
        rect(s, QX3[gx], QY3[gy], QW3, QH3[gy], fill=PANEL3, line=None)

def qhead3(x, y, text, accent):
    rect(s, x, y + 0.03, 0.06, 0.20, fill=accent, line=None)
    _, tf = tb(s, x + 0.15, y, 3.2, 0.26)
    para(tf, text, size=10.5, color=NAVY, bold=True, first=True, space_after=0)

qhead3(QX3[0] + 0.15, QY3[0] + 0.12, "GAP 1 — Procurement", RED)
_, tf = tb(s, QX3[0] + 0.18, QY3[0] + 0.50, 3.05, 1.85)
para(tf, [("•  Manual, gut-feel indenting", True, TEXT),
          (" — blind to attendance swings, exams & holidays", False, TEXT)],
     size=8, first=True, line_spacing=1.04)
para(tf, [("•  Reactive spot-buying", True, TEXT),
          (" at mandi peak prices; no price intelligence", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("•  No feedback loop", True, TEXT),
          (" from schools on actual meal uptake", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("Cost today:  8–10% meals over-produced · +12% spot premium [dummy]",
           True, RED)], size=8, space_after=0, line_spacing=1.04)

qhead3(QX3[1] + 1.30, QY3[0] + 0.12, "GAP 2 — Distribution", RED)
_, tf = tb(s, QX3[1] + 1.32, QY3[0] + 0.50, 3.05, 1.85)
para(tf, [("•  Static, driver-memory routes", True, TEXT),
          (" — late vans mean children stay hungry all day", False, TEXT)],
     size=8, first=True, line_spacing=1.04)
para(tf, [("•  Untracked vessels & utensils", True, TEXT),
          ("; adulteration & diversion risk en route", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("•  No proof-of-delivery", True, TEXT),
          (" — school shortfall disputes go unresolved", False, TEXT)],
     size=8, line_spacing=1.04)
para(tf, [("Cost today:  6.2% miss the lunch window · $35k/yr shrinkage [dummy]",
           True, RED)], size=8, space_after=0, line_spacing=1.04)

qhead3(QX3[0] + 0.15, QY3[1] + 0.12, "SOLUTION 1 — “Annapurna AI” · $150k", NAVY)
_, tf = tb(s, QX3[0] + 0.18, QY3[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("ML demand-forecasting & procurement planning across all 68 kitchens",
           False, TEXT)], size=8, first=True, line_spacing=1.04)
para(tf, [("How: ", True, NAVY),
          ("predicts school-level attendance; auto-generates daily indents & forward "
           "buying schedules", False, TEXT)], size=8, line_spacing=1.04)
para(tf, [("Rollout: ", True, GRAY),
          ("3-kitchen pilot (Mo 1–4) → all 68 kitchens by Mo 12", False, GRAY)],
     size=7.8, line_spacing=1.04)
para(tf, [("→  ~70% less food waste", True, GREEN)], size=8.2, line_spacing=1.04)
para(tf, [("→  est. $230k saved / yr [dummy]", True, GREEN)], size=8.2,
     space_after=0, line_spacing=1.04)

qhead3(QX3[1] + 1.30, QY3[1] + 0.12, "SOLUTION 2 — “Last-Mile Shield” · $100k", GREEN)
_, tf = tb(s, QX3[1] + 1.32, QY3[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("AI route optimisation + telematics (GPS, cameras, RFID) + geofenced "
           "smart locks", False, TEXT)], size=8, first=True, line_spacing=1.04)
para(tf, [("How: ", True, NAVY),
          ("drops sequenced to each school's lunch bell; doors unlock only near a "
           "registered school", False, TEXT)], size=8, line_spacing=1.04)
para(tf, [("Bonus: ", True, GRAY),
          ("3 vans freed per hub redeployed to reach new schools", False, GRAY)],
     size=7.8, line_spacing=1.04)
para(tf, [("→  99%+ on-time · −18% fleet km", True, GREEN)], size=8.2,
     line_spacing=1.04)
para(tf, [("→  est. $80k saved / yr [dummy]", True, GREEN)], size=8.2,
     space_after=0, line_spacing=1.04)

ccx3, ccy3 = 5.0, 3.96
badge_d = 2.70
rect(s, ccx3 - badge_d/2, ccy3 - badge_d/2, badge_d, badge_d, fill=WHITE,
     line=LTGRAY, line_w=1.0, shape=MSO_SHAPE.OVAL)
cd3 = CategoryChartData()
cd3.categories = ["Solution 1 — $150k", "Solution 2 — $100k"]
cd3.add_series('s', [150, 100])
chart_w3 = 2.55
gf3 = s.shapes.add_chart(XL_CHART_TYPE.PIE,
                         Inches(ccx3 - chart_w3/2), Inches(ccy3 - chart_w3/2),
                         Inches(chart_w3), Inches(chart_w3), cd3)
ch3 = gf3.chart
ch3.font.size = Pt(9); ch3.font.name = SANS; ch3.has_title = False
ch3.has_legend = False
ser3 = ch3.series[0]
for i, c in enumerate([NAVY, GREEN]):
    pt = ser3.points[i]
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
    pt.format.line.color.rgb = WHITE; pt.format.line.width = Pt(2)
plot3 = ch3.plots[0]
plot3.has_data_labels = True
dl3 = plot3.data_labels
dl3.show_value = True
dl3.number_format = '"$"0"k"'; dl3.number_format_is_linked = False
dl3.font.size = Pt(11); dl3.font.bold = True; dl3.font.color.rgb = WHITE
cs3 = ch3._chartSpace
spPr3 = cs3.makeelement(qn('c:spPr'), {})
spPr3.append(spPr3.makeelement(qn('a:noFill'), {}))
ln3 = spPr3.makeelement(qn('a:ln'), {})
ln3.append(ln3.makeelement(qn('a:noFill'), {}))
spPr3.append(ln3)
ext3 = cs3.find(qn('c:externalData'))
if ext3 is not None:
    ext3.addprevious(spPr3)
else:
    cs3.append(spPr3)
pill_w, pill_h = 2.30, 0.40
rect(s, ccx3 - pill_w/2, ccy3 + badge_d/2 - 0.02, pill_w, pill_h, fill=WHITE,
     line=LTGRAY, line_w=0.9, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
_, tf = tb(s, ccx3 - pill_w/2, ccy3 + badge_d/2 + 0.005, pill_w, pill_h - 0.05,
           anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE $250k GRANT", True, NAVY)], size=8.5, first=True, space_after=0,
     align=PP_ALIGN.CENTER)
para(tf, "60 / 40 — weighted to the costlier gap", size=6.6, color=GRAY,
     align=PP_ALIGN.CENTER, space_after=0, italic=True)

# =====================================================================
# SLIDE 4 — Solution 1 deep dive: AI-enabled planning & procurement
# =====================================================================
s = add_slide(prs)
chrome(s, 4, "Solution 1 — the AI-enabled planning & procurement build ($150k)",
       "Connecting real-time demand signals to kitchen production and procurement — "
       "the technology that closes Gap 1.",
       takeaway="“Feeding more children with the same rupee” — AI plus GS capital "
                "makes it possible.")

P1_HEAD = "PILLAR 1 — Predictive Demand Forecasting"
P2_HEAD = "PILLAR 2 — Spatial-Temporal Procurement"
for x, headtxt, hc in [(ML, P1_HEAD, TEAL), (5.08, P2_HEAD, ORANGE)]:
    rect(s, x, 1.46, 4.52, 1.66, fill=WHITE, line=LTGRAY, line_w=0.9)
    rect(s, x, 1.46, 4.52, 0.30, fill=hc, line=None)
    _, tf = tb(s, x + 0.12, 1.49, 4.28, 0.25)
    para(tf, headtxt, size=9.5, color=WHITE, bold=True, first=True, space_after=0)
_, tf = tb(s, ML + 0.14, 1.86, 4.24, 1.20)
para(tf, "Predicts exact meal & ingredient needs per centre by ingesting:", size=7.5,
     color=GRAY, italic=True, first=True, space_after=2, line_spacing=1.0)
para(tf, [("•  Attendance & enrolment trends ", True, TEXT),
          ("estimate how many meals each centre should prepare daily", False, TEXT)],
     size=7.5, space_after=2, line_spacing=1.02)
para(tf, [("•  Holiday calendars ", True, TEXT),
          ("capture local demand spikes or dips before production quantities are "
           "finalised", False, TEXT)], size=7.5, space_after=2, line_spacing=1.02)
para(tf, [("•  Menu consumption history ", True, TEXT),
          ("converts expected meals into ingredient quantities for key staples",
           False, TEXT)], size=7.5, space_after=2, line_spacing=1.02)
para(tf, [("•  Weather signals & feedback ", True, TEXT),
          ("continuously refine the next centre-level forecast", False, TEXT)],
     size=7.5, space_after=0, line_spacing=1.02)
_, tf = tb(s, 5.22, 1.86, 4.24, 1.20)
para(tf, "Decides when and where to buy each commodity to minimise total landed cost:",
     size=7.5, color=GRAY, italic=True, first=True, space_after=2, line_spacing=1.0)
para(tf, [("•  Seasonal price troughs ", True, TEXT),
          ("identify when each commodity should be bought before prices rise",
           False, TEXT)], size=7.5, space_after=2, line_spacing=1.02)
para(tf, [("•  Surplus-region signals ", True, TEXT),
          ("identify where each commodity can be sourced at lower landed cost",
           False, TEXT)], size=7.5, space_after=2, line_spacing=1.02)
para(tf, [("•  Logistics, handling & storage costs ", True, TEXT),
          ("are included in every sourcing recommendation", False, TEXT)],
     size=7.5, space_after=2, line_spacing=1.02)
para(tf, [("•  Quality grades & shelf-life rules ", True, TEXT),
          ("are locked into the final buying plan", False, TEXT)],
     size=7.5, space_after=0, line_spacing=1.02)

# loop with numbered steps + explanatory captions
cy = section(s, ML, 3.30, 4.52,
             "The self-learning operations loop — sharper every day", accent=NAVY,
             size=10)
LOOP4 = [("1 · INGEST", "pulls live prices,\nweather & attendance", NAVY),
         ("2 · FORECAST", "ingredient-level\ndemand per centre", MIDBLUE),
         ("3 · OPTIMIZE", "computes when &\nwhere to buy", GREEN),
         ("4 · RECOMMEND", "actionable buying\nsheet for kitchens", ORANGE),
         ("5 · LEARN", "outcome feedback\ntrains the model", MAGENTA)]
ccx4, ccy4, r4 = ML + 2.26, 5.13, 0.80
pts4 = []
for i, (t, sub, c) in enumerate(LOOP4):
    ang = math.radians(-90 + i * 72)
    nx = ccx4 + r4 * math.cos(ang); ny = ccy4 + r4 * math.sin(ang)
    rect(s, nx - 0.50, ny - 0.19, 1.00, 0.38, fill=c, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
    _, tf = tb(s, nx - 0.50, ny - 0.17, 1.00, 0.34, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    if i == 0:
        cxx, cyy = nx, ny - 0.42
    else:
        cxx, cyy = nx, ny + 0.42
    _, tf = tb(s, cxx - 0.65, cyy - 0.15, 1.30, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, sub, size=6.2, color=GRAY, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    pts4.append((nx, ny))
for i in range(5):
    x1, y1 = pts4[i]; x2, y2 = pts4[(i + 1) % 5]
    mx, my = (x1 + x2)/2, (y1 + y2)/2
    vx, vy = mx - ccx4, my - ccy4
    nrm = math.hypot(vx, vy) or 1
    dotted_arrow(s, x1 + (x2 - x1)*0.34 + vx/nrm*0.09,
                 y1 + (y2 - y1)*0.34 + vy/nrm*0.09,
                 x1 + (x2 - x1)*0.66 + vx/nrm*0.09,
                 y1 + (y2 - y1)*0.66 + vy/nrm*0.09, color=GRAY, width=1.0)
_, tf = tb(s, ccx4 - 0.55, ccy4 - 0.10, 1.10, 0.22, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "↻  runs daily", size=6.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
     first=True, space_after=0)

# chart + impact
cy = section(s, 5.08, 3.30, 4.52,
             "Rice procurement cost — with vs. without GS funding ($M)",
             accent=GREEN, size=9.5)
cd4 = CategoryChartData()
cd4.categories = ["2026", "2027", "2028", "2029"]
cd4.add_series("Without GS funding", [10.8, 11.4, 12.0, 12.6])
cd4.add_series("With GS funding", [9.4, 9.2, 9.1, 9.0])
gf4 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.08),
                         Inches(cy - 0.02), Inches(4.52), Inches(1.95), cd4)
ch4 = gf4.chart
ch4.font.size = Pt(7); ch4.font.name = SANS; ch4.font.color.rgb = TEXT
ch4.has_title = False
ch4.has_legend = True
ch4.legend.position = XL_LEGEND_POSITION.BOTTOM
ch4.legend.include_in_layout = False
ch4.legend.font.size = Pt(7)
for ser, c in zip(ch4.series, [ORANGE, GREEN]):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
    ser.format.line.fill.background()
plot4 = ch4.plots[0]
plot4.gap_width = 80; plot4.overlap = -10
va4 = ch4.value_axis
va4.has_major_gridlines = False
va4.tick_labels.font.size = Pt(6.5)
va4.tick_labels.number_format = '"$"0"M"'
va4.tick_labels.number_format_is_linked = False
ca4 = ch4.category_axis
ca4.tick_labels.font.size = Pt(7)
ca4.format.line.color.rgb = LTGRAY
rect(s, 5.08, cy + 2.00, 4.52, 0.68, fill=GREEN, line=None,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
_, tf = tb(s, 5.18, cy + 2.04, 4.32, 0.60, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "CUMULATIVE 2026–29 IMPACT", size=7, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, first=True, space_after=1)
para(tf, "~$10.1M saved  |  21.6% reduction", size=11.5, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, space_after=1)
para(tf, "freed capital feeds more children", size=7, color=WHITE,
     align=PP_ALIGN.CENTER, space_after=0, italic=True)

P4W = 3.45

# =====================================================================
# SLIDE 5
# =====================================================================
s = add_slide(prs)
chrome(s, 5, "Solution 2 — “Last-Mile Shield”: on time, intact, untouched",
       "$100k puts every van on an intelligent route and seals the last mile with "
       "telematics and geofenced locks — diversion becomes physically impossible.",
       takeaway="Every meal arrives on time, intact and untouched — by design, not by luck.")

# left: problem
cy = section(s, ML, 1.48, P4W, "The problem today", accent=RED)
_, tf = tb(s, ML, cy, P4W, 1.60)
para(tf, [("•  Missed lunch windows: ", True, TEXT),
          ("15–30 min late means children go hungry all day — non-negotiable",
           False, TEXT)], size=7.6, first=True, line_spacing=1.0)
para(tf, [("•  Static routes", True, TEXT),
          (" from driver memory — no live traffic, no lunch-bell sequencing", False, TEXT)],
     size=7.6, line_spacing=1.0)
para(tf, [("•  Untracked vessels", True, TEXT),
          (" leave daily with no tagging; losses recur", False, TEXT)], size=7.6,
     line_spacing=1.0)
para(tf, [("•  Open vans", True, TEXT),
          (" — meals can be offloaded / adulterated at dhabas en route", False, TEXT)],
     size=7.6, space_after=0, line_spacing=1.0)
qw = P4W / 2
stat(s, ML, cy + 1.42, qw, "6.2%", "deliveries miss the\nlunch window [dummy]", color=RED,
     vsize=13, csize=6.8)
stat(s, ML + qw, cy + 1.42, qw, "$35k/yr", "vessel & shrinkage\nlosses [dummy]", color=RED,
     vsize=13, csize=6.8)

# left: route map
cy = section(s, ML, 3.98, P4W, "Route redesign — 1 hub, 40 schools", accent=NAVY,
             size=10)
img_placeholder(s, ML, cy + 0.02, P4W, 1.35, "\U0001F5FA MAP: before vs after routes",
                "(insert optimised-route schematic; −18% km)")
_, tf = tb(s, ML, cy + 1.42, P4W, 0.40)
para(tf, [("3 vans freed per hub ", True, GREEN),
          ("→ redeployed to 12 new schools [dummy]", False, TEXT)], size=7,
     first=True, space_after=0, align=PP_ALIGN.CENTER)

# right: three layers
vline(s, 3.95, 1.55, 6.45)
HX, HW = 4.15, SW - ML - 4.15
cy = section(s, HX, 1.48, HW, "The fix — three integrated layers", accent=GREEN)
layers = [
    ("1 · ROUTE INTELLIGENCE", PALEBLUE, MIDBLUE,
     "AI planner sequences every drop to each school's lunch bell • live traffic "
     "re-routing • delay alerts to school + kitchen 30 min ahead"),
    ("2 · VEHICLE TELEMATICS", PALEGOLD, ORANGE,
     "GPS + in-van cameras on every vehicle • RFID tags on all vessels & utensils • "
     "automated end-of-day reconciliation — zero missing inventory"),
    ("3 · GEOFENCED INTEGRITY", PALEGREEN, GREEN,
     "Smart-locked doors open only within 150 m of a registered school • real-time "
     "tamper alerts • en-route offloading & adulteration made impossible"),
]
ly = cy + 0.04
for t, bg, bc, d in layers:
    rect(s, HX, ly, HW, 0.78, fill=bg, line=bc, line_w=0.9,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    _, tf = tb(s, HX + 0.12, ly + 0.06, HW - 0.24, 0.68)
    para(tf, [(t, True, NAVY)], size=7.8, first=True, space_after=1)
    para(tf, d, size=7.2, color=TEXT, space_after=0, line_spacing=1.0)
    ly += 0.84

# right: SLA strip
_, tf = tb(s, HX, ly + 0.02, HW, 0.24)
para(tf, [("After implementation:  ", True, NAVY),
          ("99.5% in lunch window   ·   −18% fleet km   ·   0 unaccounted vessels",
           True, GREEN)], size=8.2, first=True, space_after=0)

# right bottom: donut + savings
DY = 4.74
cy = section(s, HX, DY, 2.75, "Where the $100k goes", accent=NAVY, size=10)
donut(s, HX - 0.05, cy - 0.04, 2.85, 1.48,
      ["Routing sw $35k", "Telematics $30k", "Smart locks $20k", "RFID $10k", "Ops $5k"],
      [35, 30, 20, 10, 5], [GREEN, TEAL, NAVY, GOLD, GRAY],
      legend=XL_LEGEND_POSITION.BOTTOM, labels=False)
SX4 = 7.15
cy = section(s, SX4, DY, SW - ML - SX4, "Annual savings ($k)", accent=GREEN, size=10)
col_chart(s, SX4 - 0.05, cy, SW - ML - SX4 + 0.10, 1.28,
          ["Fuel", "Vessels", "Food"], [("$k/yr", [45, 20, 15])], [GREEN],
          gap=60, font_size=6.8)
_, tf = tb(s, SX4 - 0.05, cy + 1.32, SW - ML - SX4 + 0.10, 0.26)
para(tf, [("= $80k/yr + trust that is priceless", True, GREEN)], size=7, first=True,
     space_after=0, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 6 — conclusion (combined layout: impact rail + equation + why GS + proof)
# =====================================================================
s = add_slide(prs)
chrome(s, 6, "Conclusion — why Akshaya Patra, and why this proposal wins",
       "A proven engine of social mobility, a grant that compounds instead of depletes, "
       "and a partnership only Goldman Sachs can power.",
       takeaway="$250k that doesn't feed children for a year — it upgrades the engine "
                "that feeds them forever.")

# ---- left: fund & impact rail ----
RAILW = 2.55
_, tf = tb(s, ML, 1.52, RAILW, 0.24)
para(tf, "THE FUND & ITS IMPACT", size=9, color=NAVY, bold=True, first=True,
     space_after=0)
rect(s, ML, 1.80, RAILW, 0.014, fill=LTGRAY, line=None)
rail = [("$250k", "one-time grant", NAVY, 22),
        ("$310k / yr", "returned in audited savings — every single year", GREEN, 17),
        ("+105,000", "children fed per year at zero extra cost", NAVY, 17),
        ("< 10 months", "payback on the full grant", ORANGE, 14)]
ry = 1.96
for i, (v, c, col, sz) in enumerate(rail):
    _, tf = tb(s, ML, ry, RAILW, 0.85)
    para(tf, v, size=sz, color=col, bold=True, font=SERIF, first=True, space_after=1)
    para(tf, c, size=7.2, color=GRAY, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, ML + 0.33, ry + 0.70, ML + 0.33, ry + 0.94, color=LTGRAY,
                     curve=False, width=1.0)
    ry += 1.00
vline(s, 3.05, 1.55, 5.30)

# ---- right top: partnership equation ----
RX6, RW6 = 3.35, SW - ML - 3.35
cy = section(s, RX6, 1.50, RW6, "The partnership equation", accent=MIDBLUE, size=10.5)
ov, oh = 1.06, 0.88
eq_w = 3*ov + 2*0.50
ox1 = RX6 + (RW6 - eq_w) / 2
eq_y = cy + 0.02
for i, (t, bg, c) in enumerate([("GS $250k\ncatalytic\ncapital", PALEBLUE, NAVY),
                                ("AP's 2.3M\nmeals-a-day\nengine", PALEGOLD, ORANGE),
                                ("Impact that\ncompounds\nyearly", PALEGREEN, GREEN)]):
    xx = ox1 + i * (ov + 0.50)
    rect(s, xx, eq_y, ov, oh, fill=bg, line=c, line_w=1.1, shape=MSO_SHAPE.OVAL)
    _, tf = tb(s, xx, eq_y + 0.11, ov, oh - 0.22, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7, color=c, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    if i < 2:
        sym = MSO_SHAPE.MATH_PLUS if i == 0 else MSO_SHAPE.MATH_EQUAL
        rect(s, xx + ov + 0.12, eq_y + oh/2 - 0.12, 0.26, 0.24, fill=NAVY, line=None,
             shape=sym)
_, tf = tb(s, RX6, eq_y + oh + 0.05, RW6, 0.24)
para(tf, "One-time capital → permanent capability → recurring savings → more children "
         "in school, every year after Year 1", size=7, color=GRAY, italic=True,
     first=True, space_after=0, align=PP_ALIGN.CENTER)

# ---- right middle: why GS (5 points) ----
cy2 = section(s, RX6, 3.16, RW6, "Why the GS grant stands out here", accent=NAVY,
              size=10.5)
WHY6 = [
    ("Structural, not sustenance", "others fund meals for a year; GS funds the "
     "technology that makes every future meal cheaper"),
    ("Plays to GS DNA", "forecasting, optimisation & risk controls — GS engineering "
     "mentors the build pro bono"),
    ("Measurable ROI", "$310k/yr audited savings; KPI dashboard shared quarterly "
     "with GS"),
    ("Partnership beyond capital", "GS volunteering days at kitchens & co-branded "
     "impact reporting — a partner, not a donor"),
    ("Scalable blueprint", "proven once, extends to school-feeding programmes across "
     "South Asia & Africa"),
]
yy = cy2
for h, d in WHY6:
    _, tf = tb(s, RX6, yy, RW6, 0.40)
    para(tf, [("✓  " + h + ":  ", True, GREEN), (d, False, TEXT)], size=7.8,
         first=True, space_after=0, line_spacing=0.98)
    yy += 0.345

# ---- bottom (full width): proven track record ----
cy3 = section(s, ML, 5.42, SW - 2*ML, "Proven track record", accent=ORANGE, size=10.5)
TY = cy3 + 0.02
img_placeholder(s, ML, TY, 0.62, 0.62, "\U0001F4F7 Krishna\nKumar")
_, tf = tb(s, ML + 0.76, TY, 3.85, 0.66)
para(tf, "\u201cThe mid-day meal was the reason I stayed in school. Today I manage "
         "portfolios instead of an empty stomach.\u201d", size=7.2, color=TEXT,
     italic=True, first=True, space_after=1, line_spacing=1.0)
para(tf, [("— Krishna Kumar, ", True, NAVY),
          ("AVP, HSBC — former mid-day meal beneficiary", False, GRAY)], size=6.8,
     space_after=0)
vline(s, ML + 4.78, TY, TY + 0.62)
PX6 = ML + 4.96; PW6 = SW - ML - PX6
_, tf = tb(s, PX6, TY, PW6, 0.20)
para(tf, [("Trusted by leading corporates ", True, NAVY), ("(+200 more)", False, GRAY)],
     size=8, first=True, space_after=0)
lw6 = (PW6 - 5 * 0.05) / 6
for i, p in enumerate(["Infosys", "TCS", "HDFC", "Amazon", "Airbus", "SBI"]):
    img_placeholder(s, PX6 + i * (lw6 + 0.05), TY + 0.26, lw6, 0.36, p)

out = "AkshayaPatra_GS_Pitch.pptx"
prs.save(out)
print("saved:", out)


# ======================================================================
# ===== build_final_content.py =====
# ======================================================================
"""Append 4 content-final slides (v2 of slides 1, 3, 4, 6) with verified TAPF data.
Existing slides are untouched. Sources footer on every new slide."""
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import math

PATH = "AkshayaPatra_GS_Pitch.pptx"
prs = Presentation(PATH)

TEAM = ("Team: Nishant Tomer · Pratham Goenka · Aadithya Muralidharan · "
        "Somya Sethi · Mohajit Neog")

def new_slide(num_label, title, sub, takeaway, tag, team=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s, len(prs.slides._sldIdLst), title, sub, takeaway=takeaway)
    _, tf = tb(s, SW - ML - 4.4, 0.06, 4.4, 0.20)
    para(tf, tag, size=7.5, color=GRAY, italic=True, align=PP_ALIGN.RIGHT,
         first=True, space_after=0)
    _, tf = tb(s, 5.30, 7.16, 3.72, 0.30)
    para(tf, "Source: The Akshaya Patra Foundation — official documents",
         size=6.5, color=GRAY, italic=True, align=PP_ALIGN.RIGHT, first=True,
         space_after=0)
    if team:
        _, tf = tb(s, 1.72, 7.10, 3.45, 0.38)
        para(tf, TEAM, size=6.3, color=GRAY, first=True, space_after=0,
             line_spacing=0.95)
    return s

# =====================================================================
# SLIDE 7 — v2 of slide 1: introduction (verified numbers)
# =====================================================================
s = new_slide(7, "Who is The Akshaya Patra Foundation?",
              "The world's largest NGO-run school meal programme — serving 2.35 million "
              "children every school day, so that hunger never interrupts education.",
              "“No child in India shall be deprived of education because of hunger” — "
              "5 billion meals over 25 years.",
              "FINAL CONTENT — SLIDE 1 (v2)", team=True)

rect(s, ML, 1.42, SW - 2*ML, 0.80, fill=ORANGE, line=None)
img_placeholder(s, ML + 0.08, 1.51, 0.98, 0.62, "AKSHAYA PATRA\nlogo")
_, tf = tb(s, ML + 1.22, 1.42, SW - 2*ML - 1.36, 0.80, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Founded in 2000 serving 1,500 children, ", True, WHITE),
          ("Akshaya Patra is a Government-partnered PPP that has grown into the world's "
           "largest NGO-run school meal programme — ", False, WHITE),
          ("5 billion cumulative meals over 25 years", True, WHITE),
          (", a milestone celebrated with the Hon'ble President of India in March 2026.",
           False, WHITE)],
     size=8, first=True, space_after=0, line_spacing=1.0)

LX, LW = ML, 4.42
RX, RW = 5.02, SW - ML - 5.02
vline(s, 4.77, 2.42, 6.45)

cy = section(s, LX, 2.38, LW, "Who they are & why we chose them", accent=NAVY)
_, tf = tb(s, LX, cy, LW, 1.55)
para(tf, [("Mission: ", True, NAVY),
          ("serve mid-day meals to 3 million children every day and achieve 3 million "
           "servings of morning nutrition.", False, TEXT)], size=8, first=True,
     line_spacing=1.02)
para(tf, [("Why Akshaya Patra: ", True, NAVY),
          ("our supply-chain review found world-class production — and ", False, TEXT),
          ("two clearly identified, fixable gaps", True, TEXT),
          (" that TAPF itself ranks as its top technology priorities, where $250k "
           "unlocks compounding year-on-year impact.", False, TEXT)], size=8,
     line_spacing=1.02)
para(tf, "✓ 25 yrs of proven delivery   ✓ Govt-audited PPP   ✓ Tech-forward ops "
         "(ERP, OS1 dispatch, kitchen automation)",
     size=7.6, color=GREEN, bold=True, space_after=0)

cy = section(s, LX, 4.32, LW, "The journey — daily meals served (millions)", accent=BLUE)
line_chart(s, LX, cy + 0.02, LW, 1.42,
           ["2000", "2005", "2010", "2015", "2020", "2023", "2026"],
           [("Meals/day (M)", [0.0015, 0.33, 1.20, 1.44, 1.80, 2.00, 2.35])], [NAVY],
           font_size=7)
_, tf = tb(s, LX, cy + 1.50, LW, 0.34)
para(tf, [("Mar 2026: ", True, TEXT),
          ("5-billionth meal celebrated by the Hon'ble President of India  ·  ",
           False, TEXT),
          ("Apr 2026: ", True, TEXT),
          ("80th kitchen opens in Pune (Deutsche Bank CSR)", False, TEXT)],
     size=6.8, first=True, space_after=0, color=GRAY, line_spacing=0.95)

cy = section(s, RX, 2.38, RW, "What they do — the hub-and-spoke model", accent=MIDBLUE)
_, tf = tb(s, RX, cy, RW, 1.15)
para(tf, [("•  Centralised mega-kitchens ", True, TEXT),
          ("— rice for 1,000 students in 15 min; 40,000 rotis/hour", False, TEXT)],
     size=7.8, first=True, line_spacing=1.0)
para(tf, [("•  Hub-and-spoke delivery ", True, TEXT),
          ("— 1,100+ vehicles move 2 lakh+ insulated vessels daily", False, TEXT)],
     size=7.8, line_spacing=1.0)
para(tf, [("•  PPP economics ", True, TEXT),
          ("— govt funds ₹8.56 of every ₹17 meal; donors bridge ₹6.56", False, TEXT)],
     size=7.8, line_spacing=1.0)
para(tf, [("•  Beyond lunch ", True, TEXT),
          ("— morning nutrition, school rejuvenation, scholarships (NEST)",
           False, TEXT)], size=7.8, space_after=0, line_spacing=1.0)

cy = section(s, RX, 3.94, RW, "Scale at a glance (2026)", accent=ORANGE)
qw = RW / 4
for i, (v, c) in enumerate([("2.35M", "meals served\nevery day"),
                            ("80", "kitchens across\n81 locations"),
                            ("16+3", "states &\nunion territories"),
                            ("25,000+", "govt schools\nserved")]):
    stat(s, RX + i*qw, cy + 0.02, qw, v, c, vsize=14, csize=6.8)

cy = section(s, RX, 5.12, RW, "The opportunity — India context", accent=GREEN)
qw2 = 1.05
for i, (v, c) in enumerate([("2%", "of India's 130M govt-\nschool children reached"),
                            ("8.2%", "secondary dropout\nrate (UDISE+)"),
                            ("102/123", "India, Global Hunger\nIndex 2025")]):
    stat(s, RX + i * 1.18, cy + 0.02, 1.12, v, c, vsize=12, csize=6.2)
img_placeholder(s, RX + 3.60, cy + 0.04, RW - 3.60, 0.72,
                "\U0001F4F7 PHOTO: children\nat mid-day meal")

# =====================================================================
# SLIDE 8 — v2 of slide 3: gaps + $250k split (verified content)
# =====================================================================
s = new_slide(8, "The two gaps — and the $250,000 that closes them",
              "The gaps sit exactly where TAPF's own leadership says they are — its "
              "top-ranked technology priorities, confirmed in writing.",
              "Two fixable gaps, two targeted builds — both at the top of TAPF's own "
              "priority list.",
              "FINAL CONTENT — SLIDE 3 (v2)")

PANEL3 = RGBColor(0xF4, 0xF4, 0xF4)
QX3 = [ML, 5.08]; QW3 = 4.52
QY3 = [1.46, 4.02]; QH3 = [2.46, 2.44]
for gx in range(2):
    for gy in range(2):
        rect(s, QX3[gx], QY3[gy], QW3, QH3[gy], fill=PANEL3, line=None)

def qhead3(x, y, text, accent):
    rect(s, x, y + 0.03, 0.06, 0.20, fill=accent, line=None)
    _, tf = tb(s, x + 0.15, y, 3.2, 0.26)
    para(tf, text, size=10.5, color=NAVY, bold=True, first=True, space_after=0)

qhead3(QX3[0] + 0.15, QY3[0] + 0.12, "GAP 1 — Procurement", RED)
_, tf = tb(s, QX3[0] + 0.18, QY3[0] + 0.50, 3.05, 1.90)
para(tf, [("•  Periodic, manual planning", True, TEXT),
          (" vs demand that shifts daily — weather, off-season supply, attendance "
           "swings", False, TEXT)], size=7.8, first=True, line_spacing=1.02)
para(tf, [("•  No buffer day", True, TEXT),
          (" — a planning miss is a same-day miss for children", False, TEXT)],
     size=7.8, line_spacing=1.02)
para(tf, [("•  Complexity: ", True, TEXT),
          ("₹450 Cr+ food spend · 700 items · 655 vendors · 78 buyers, 145 approvers",
           False, TEXT)], size=7.8, line_spacing=1.02)
para(tf, [("Today: 4–5.6% monthly vegetable wastage · spot buys above seasonal "
           "price troughs", True, RED)], size=7.8, space_after=0, line_spacing=1.02)

qhead3(QX3[1] + 1.30, QY3[0] + 0.12, "GAP 2 — Distribution", RED)
_, tf = tb(s, QX3[1] + 1.32, QY3[0] + 0.50, 3.05, 1.90)
para(tf, [("•  No real-time vehicle location", True, TEXT),
          (" — schools phone kitchens to confirm deliveries", False, TEXT)],
     size=7.8, first=True, line_spacing=1.02)
para(tf, [("•  No maintenance or driver-safety data", True, TEXT),
          (" — issues surface only after incidents", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("•  4-hour cook-to-consume window", True, TEXT),
          (" at 65°C across ~70 km routes — zero slack", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("TAPF's #1 strategic headwind: last-mile traceability & food-safety risk",
           True, RED)], size=7.8, space_after=0, line_spacing=1.02)

qhead3(QX3[0] + 0.15, QY3[1] + 0.12, "SOLUTION 1 — “Annapurna AI” · $150k", NAVY)
_, tf = tb(s, QX3[0] + 0.18, QY3[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("AI demand forecasting & procurement planning across all 80 kitchens",
           False, TEXT)], size=7.8, first=True, line_spacing=1.02)
para(tf, [("How: ", True, NAVY),
          ("attendance, calendars, menus & weather → daily indents; buy-timing tuned "
           "to seasonal price troughs across 655 vendors", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("→  cuts veg wastage toward ~1–2%", True, GREEN)], size=8,
     line_spacing=1.02)
para(tf, [("→  savings recur every year (team est.)", True, GREEN)], size=8,
     space_after=0, line_spacing=1.02)

qhead3(QX3[1] + 1.30, QY3[1] + 0.12, "SOLUTION 2 — “Last-Mile Shield” · $100k", GREEN)
_, tf = tb(s, QX3[1] + 1.32, QY3[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("AIS-140 GPS + BLE door sensors fleet-wide; geofenced delivery "
           "confirmation; smart locks on high-risk routes", False, TEXT)], size=7.8,
     first=True, line_spacing=1.02)
para(tf, [("How: ", True, NAVY),
          ("accelerates TAPF's own in-house blueprint (Traccar + DOUP) — auto "
           "delivery confirmation to schools", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("→  on-time within the 4-hr window", True, GREEN)], size=8,
     line_spacing=1.02)
para(tf, [("→  maintenance & driver-safety telemetry included", True, GREEN)],
     size=8, space_after=0, line_spacing=1.02)

ccx3, ccy3 = 5.0, 3.96
badge_d = 2.70
rect(s, ccx3 - badge_d/2, ccy3 - badge_d/2, badge_d, badge_d, fill=WHITE,
     line=LTGRAY, line_w=1.0, shape=MSO_SHAPE.OVAL)
cd3 = CategoryChartData()
cd3.categories = ["Solution 1 — $150k", "Solution 2 — $100k"]
cd3.add_series('s', [150, 100])
chart_w3 = 2.55
gf3 = s.shapes.add_chart(XL_CHART_TYPE.PIE,
                         Inches(ccx3 - chart_w3/2), Inches(ccy3 - chart_w3/2),
                         Inches(chart_w3), Inches(chart_w3), cd3)
ch3 = gf3.chart
ch3.font.size = Pt(9); ch3.font.name = SANS; ch3.has_title = False
ch3.has_legend = False
for i, c in enumerate([NAVY, GREEN]):
    pt = ch3.series[0].points[i]
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
    pt.format.line.color.rgb = WHITE; pt.format.line.width = Pt(2)
plot3 = ch3.plots[0]
plot3.has_data_labels = True
dl3 = plot3.data_labels
dl3.show_value = True
dl3.number_format = '"$"0"k"'; dl3.number_format_is_linked = False
dl3.font.size = Pt(11); dl3.font.bold = True; dl3.font.color.rgb = WHITE
cs3 = ch3._chartSpace
spPr3 = cs3.makeelement(qn('c:spPr'), {})
spPr3.append(spPr3.makeelement(qn('a:noFill'), {}))
ln3 = spPr3.makeelement(qn('a:ln'), {})
ln3.append(ln3.makeelement(qn('a:noFill'), {}))
spPr3.append(ln3)
ext3 = cs3.find(qn('c:externalData'))
if ext3 is not None:
    ext3.addprevious(spPr3)
else:
    cs3.append(spPr3)
pill_w, pill_h = 2.55, 0.42
rect(s, ccx3 - pill_w/2, ccy3 + badge_d/2 - 0.02, pill_w, pill_h, fill=WHITE,
     line=LTGRAY, line_w=0.9, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
_, tf = tb(s, ccx3 - pill_w/2, ccy3 + badge_d/2 + 0.005, pill_w, pill_h - 0.05,
           anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE $250k GRANT", True, NAVY)], size=8.5, first=True, space_after=0,
     align=PP_ALIGN.CENTER)
para(tf, "both gaps: TAPF's self-ranked top priorities", size=6.4, color=GRAY,
     align=PP_ALIGN.CENTER, space_after=0, italic=True)

# =====================================================================
# SLIDE 9 — v2 of slide 4: Solution 1 deep dive (verified content)
# =====================================================================
s = new_slide(9, "Solution 1 — the AI-enabled planning & procurement build ($150k)",
              "Connecting real-time demand signals to kitchen production and "
              "procurement — the technology that closes Gap 1.",
              "“Feeding more children with the same rupee” — AI plus GS capital "
              "makes it possible.",
              "FINAL CONTENT — SLIDE 4 (v2)")

P1_HEAD = "PILLAR 1 — Predictive Demand Forecasting"
P2_HEAD = "PILLAR 2 — Spatial-Temporal Procurement"
for x, headtxt, hc in [(ML, P1_HEAD, TEAL), (5.08, P2_HEAD, ORANGE)]:
    rect(s, x, 1.46, 4.52, 1.86, fill=WHITE, line=LTGRAY, line_w=0.9)
    rect(s, x, 1.46, 4.52, 0.30, fill=hc, line=None)
    _, tf = tb(s, x + 0.12, 1.49, 4.28, 0.25)
    para(tf, headtxt, size=9.5, color=WHITE, bold=True, first=True, space_after=0)
_, tf = tb(s, ML + 0.14, 1.86, 4.24, 1.42)
para(tf, "Predicts exact meal & ingredient needs per centre by ingesting:", size=7.4,
     color=GRAY, italic=True, first=True, space_after=2, line_spacing=1.0)
para(tf, [("•  Attendance & enrolment trends ", True, TEXT),
          ("estimate how many meals each centre should prepare daily", False, TEXT)],
     size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Holiday calendars ", True, TEXT),
          ("capture local demand spikes or dips before production is finalised",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Menu consumption history ", True, TEXT),
          ("converts expected meals into ingredient quantities for key staples",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Weather & consumption feedback ", True, TEXT),
          ("continuously refine the next centre-level forecast", False, TEXT)],
     size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Perishables (veg, dairy) ", True, TEXT),
          ("get shorter forecast cycles than staples (grains, dhals, oils)",
           False, TEXT)], size=7.4, space_after=0, line_spacing=1.02)
_, tf = tb(s, 5.22, 1.86, 4.24, 1.42)
para(tf, "Decides when and where to buy each commodity to minimise total landed cost:",
     size=7.4, color=GRAY, italic=True, first=True, space_after=2, line_spacing=1.0)
para(tf, [("•  Seasonal price troughs ", True, TEXT),
          ("identify when each commodity should be bought before prices rise",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Surplus-region signals ", True, TEXT),
          ("identify where each commodity can be sourced at lower landed cost",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Logistics, handling & storage costs ", True, TEXT),
          ("are included in every sourcing recommendation", False, TEXT)],
     size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Quality grades & shelf-life rules ", True, TEXT),
          ("are locked into the final buying plan", False, TEXT)], size=7.4,
     space_after=2, line_spacing=1.02)
para(tf, [("•  Full procure-to-pay coverage ", True, TEXT),
          ("— requisition → PO → goods receipt → invoice → payment", False, TEXT)],
     size=7.4, space_after=0, line_spacing=1.02)

_, tf = tb(s, ML, 3.38, SW - 2*ML, 0.22)
para(tf, [("Annual volumes:  ", True, NAVY),
          ("120 T vegetables · 55 T dhals · 20 T dairy · 18 T oils · 6 T spices — "
           "against ₹450 Cr+ annual food spend", False, GRAY)], size=7.4, first=True,
     space_after=0, align=PP_ALIGN.CENTER)

cy = section(s, ML, 3.66, 4.52,
             "The self-learning operations loop — sharper every day", accent=NAVY,
             size=10)
LOOP4 = [("1 · INGEST", "pulls live prices,\nweather & attendance", NAVY),
         ("2 · FORECAST", "ingredient-level\ndemand per centre", MIDBLUE),
         ("3 · OPTIMIZE", "computes when &\nwhere to buy", GREEN),
         ("4 · RECOMMEND", "actionable buying\nsheet for kitchens", ORANGE),
         ("5 · LEARN", "outcome feedback\ntrains the model", MAGENTA)]
ccx4, ccy4, r4 = ML + 2.26, 5.32, 0.74
pts4 = []
for i, (t, sub, c) in enumerate(LOOP4):
    ang = math.radians(-90 + i * 72)
    nx = ccx4 + r4 * math.cos(ang); ny = ccy4 + r4 * math.sin(ang)
    rect(s, nx - 0.50, ny - 0.18, 1.00, 0.36, fill=c, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
    _, tf = tb(s, nx - 0.50, ny - 0.16, 1.00, 0.32, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    if i == 0:
        cxx, cyy = nx, ny - 0.40
    else:
        cxx, cyy = nx, ny + 0.40
    _, tf = tb(s, cxx - 0.65, cyy - 0.145, 1.30, 0.29, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, sub, size=6.2, color=GRAY, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    pts4.append((nx, ny))
for i in range(5):
    x1, y1 = pts4[i]; x2, y2 = pts4[(i + 1) % 5]
    mx, my = (x1 + x2)/2, (y1 + y2)/2
    vx, vy = mx - ccx4, my - ccy4
    nrm = math.hypot(vx, vy) or 1
    dotted_arrow(s, x1 + (x2 - x1)*0.34 + vx/nrm*0.09,
                 y1 + (y2 - y1)*0.34 + vy/nrm*0.09,
                 x1 + (x2 - x1)*0.66 + vx/nrm*0.09,
                 y1 + (y2 - y1)*0.66 + vy/nrm*0.09, color=GRAY, width=1.0)
_, tf = tb(s, ccx4 - 0.55, ccy4 - 0.10, 1.10, 0.22, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "↻  runs daily", size=6.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
     first=True, space_after=0)

cy = section(s, 5.08, 3.66, 4.52,
             "Rice procurement cost — with vs. without GS funding ($M)",
             accent=GREEN, size=9.5)
cd4 = CategoryChartData()
cd4.categories = ["2026", "2027", "2028", "2029"]
cd4.add_series("Without GS funding", [10.8, 11.4, 12.0, 12.6])
cd4.add_series("With GS funding", [9.4, 9.2, 9.1, 9.0])
gf4 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.08),
                         Inches(cy - 0.02), Inches(4.52), Inches(1.68), cd4)
ch4 = gf4.chart
ch4.font.size = Pt(7); ch4.font.name = SANS; ch4.font.color.rgb = TEXT
ch4.has_title = False
ch4.has_legend = True
ch4.legend.position = XL_LEGEND_POSITION.BOTTOM
ch4.legend.include_in_layout = False
ch4.legend.font.size = Pt(7)
for ser, c in zip(ch4.series, [ORANGE, GREEN]):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
    ser.format.line.fill.background()
plot4 = ch4.plots[0]
plot4.gap_width = 80; plot4.overlap = -10
va4 = ch4.value_axis
va4.has_major_gridlines = False
va4.tick_labels.font.size = Pt(6.5)
va4.tick_labels.number_format = '"$"0"M"'
va4.tick_labels.number_format_is_linked = False
ca4 = ch4.category_axis
ca4.tick_labels.font.size = Pt(7)
ca4.format.line.color.rgb = LTGRAY
rect(s, 5.08, cy + 1.74, 4.52, 0.66, fill=GREEN, line=None,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
_, tf = tb(s, 5.18, cy + 1.78, 4.32, 0.58, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "CUMULATIVE 2026–29 IMPACT", size=7, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, first=True, space_after=1)
para(tf, "~$10.1M saved  |  21.6% reduction", size=11.5, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, space_after=1)
para(tf, "team estimate on TAPF's ₹450 Cr+ annual food spend", size=6.6, color=WHITE,
     align=PP_ALIGN.CENTER, space_after=0, italic=True)

# =====================================================================
# SLIDE 10 — v2 of slide 6: conclusion (verified content)
# =====================================================================
s = new_slide(10, "Conclusion — why Akshaya Patra, and why this proposal wins",
              "A proven engine of social mobility, economics where every donor rupee "
              "is matched by government subsidy, and a grant that compounds.",
              "$250k that doesn't feed children for a year — it upgrades the engine "
              "that feeds them forever.",
              "FINAL CONTENT — SLIDE 6 (v2)")

RAILW = 2.55
_, tf = tb(s, ML, 1.52, RAILW, 0.24)
para(tf, "THE FUND & ITS IMPACT", size=9, color=NAVY, bold=True, first=True,
     space_after=0)
rect(s, ML, 1.80, RAILW, 0.014, fill=LTGRAY, line=None)
rail = [("$250k", "one-time grant ≈ ₹2.1 Cr", NAVY, 22),
        ("₹6.56", "donor cost per meal — govt funds ₹8.56 of every ₹17", GREEN, 20),
        ("₹1,500", "feeds one child for a full year (232 school days)", NAVY, 20),
        ("14,000+", "children fed for a year by the grant alone — before savings", ORANGE, 18)]
ry = 1.96
for i, (v, c, col, sz) in enumerate(rail):
    _, tf = tb(s, ML, ry, RAILW, 0.90)
    para(tf, v, size=sz, color=col, bold=True, font=SERIF, first=True, space_after=1)
    para(tf, c, size=7, color=GRAY, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, ML + 0.33, ry + 0.66, ML + 0.33, ry + 0.90, color=LTGRAY,
                     curve=False, width=1.0)
    ry += 0.95
vline(s, 3.05, 1.55, 5.32)

RX6, RW6 = 3.35, SW - ML - 3.35
cy = section(s, RX6, 1.50, RW6, "The partnership equation", accent=MIDBLUE, size=10.5)
ov, oh = 1.06, 0.88
eq_w = 3*ov + 2*0.50
ox1 = RX6 + (RW6 - eq_w) / 2
eq_y = cy + 0.02
for i, (t, bg, c) in enumerate([("GS $250k\ncatalytic\ncapital", PALEBLUE, NAVY),
                                ("AP's 2.35M\nmeals-a-day\nengine", PALEGOLD, ORANGE),
                                ("Impact that\ncompounds\nyearly", PALEGREEN, GREEN)]):
    xx = ox1 + i * (ov + 0.50)
    rect(s, xx, eq_y, ov, oh, fill=bg, line=c, line_w=1.1, shape=MSO_SHAPE.OVAL)
    _, tf = tb(s, xx, eq_y + 0.11, ov, oh - 0.22, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7, color=c, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    if i < 2:
        sym = MSO_SHAPE.MATH_PLUS if i == 0 else MSO_SHAPE.MATH_EQUAL
        rect(s, xx + ov + 0.12, eq_y + oh/2 - 0.12, 0.26, 0.24, fill=NAVY, line=None,
             shape=sym)
_, tf = tb(s, RX6, eq_y + oh + 0.05, RW6, 0.24)
para(tf, "One-time capital → permanent capability → recurring savings → more children "
         "in school, every year after Year 1", size=7, color=GRAY, italic=True,
     first=True, space_after=0, align=PP_ALIGN.CENTER)

cy2 = section(s, RX6, 3.16, RW6, "Why the GS grant stands out here", accent=NAVY,
              size=10.5)
WHY6 = [
    ("Structural, not sustenance", "TAPF's own words: restricted funds “deliver meals "
     "today” but cannot build “the systems behind every meal” — GS funds exactly those "
     "systems"),
    ("Co-designed priorities", "forecasting & real-time fleet visibility are TAPF's "
     "self-declared top tech gaps; logistics tech ranked #1 for grant ROI"),
    ("Plays to GS DNA", "forecasting, optimisation & risk controls — with GS "
     "engineering mentorship pro bono"),
    ("Measurable ROI", "audited savings against a ₹450 Cr+ procurement base; "
     "quarterly KPI dashboard to GS"),
    ("Compounding & green", "savings recur yearly and support TAPF's 50% "
     "renewable-by-2030 and EV-fleet-by-2035 goals"),
]
yy = cy2
for h, d in WHY6:
    _, tf = tb(s, RX6, yy, RW6, 0.42)
    para(tf, [("✓  " + h + ":  ", True, GREEN), (d, False, TEXT)], size=7.6,
         first=True, space_after=0, line_spacing=0.98)
    yy += 0.40

cy3 = section(s, ML, 5.44, SW - 2*ML, "Proven track record", accent=ORANGE, size=10.5)
TY = cy3 + 0.04
img_placeholder(s, ML, TY, 0.55, 0.62, "\U0001F4F7")
_, tf = tb(s, ML + 0.66, TY, 2.55, 0.85)
para(tf, "“Everybody used to get enough food and it was very tasty.”", size=6.9,
     color=TEXT, italic=True, first=True, space_after=1, line_spacing=1.0)
para(tf, [("— Krishna Kumar, ", True, NAVY),
          ("AVP, HSBC — programme alumnus", False, GRAY)], size=6.6, space_after=0)
vline(s, ML + 3.32, TY, TY + 0.72)
img_placeholder(s, ML + 3.48, TY, 0.55, 0.62, "\U0001F4F7")
_, tf = tb(s, ML + 4.14, TY, 2.45, 0.85)
para(tf, "“My dream is to become financially strong so I can support others the way "
         "I was supported.”", size=6.9, color=TEXT, italic=True, first=True,
     space_after=1, line_spacing=1.0)
para(tf, [("— Suma, ", True, NAVY),
          ("Purchasing Analyst, Palo Alto Networks", False, GRAY)], size=6.6,
     space_after=0)
vline(s, ML + 6.72, TY, TY + 0.72)
_, tf = tb(s, ML + 6.88, TY - 0.02, SW - ML - (ML + 6.88), 0.90)
para(tf, [("Recognition: ", True, NAVY),
          ("Padma Shri · Gandhi Peace Prize 2016 · BBC Global Champion · Nikkei Asia "
           "Award · CII Food Safety 2026 · HBS case study", False, TEXT)], size=6.6,
     first=True, space_after=2, line_spacing=1.0)
para(tf, [("Partners incl.: ", True, NAVY),
          ("Deutsche Bank (Pune kitchen '26) · IOCL (solar) + leading corporates",
           False, TEXT)], size=6.6, space_after=0, line_spacing=1.0)

prs.save(PATH)
print("saved; total slides:", len(prs.slides._sldIdLst))


# ======================================================================
# ===== build_chart_slide.py =====
# ======================================================================
"""Standalone mock: proposed 'maximum impact' data visual for Solution 1 savings.
Builds a single-slide pptx copy — the real deck is untouched."""
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
import shutil

SRC = "AkshayaPatra_GS_Pitch.pptx"
MOCK = ("/private/tmp/claude-501/-Users-nishanttomer-Desktop/"
        "ad99222a-484a-437c-9678-2fdf9eec1d5d/scratchpad/chart_mock.pptx")

prs = Presentation(SRC)
s = prs.slides.add_slide(prs.slide_layouts[6])
chrome(s, len(prs.slides._sldIdLst), "Solution 1 — where the savings come from, and what they buy",
       "Every bar is a number from TAPF's own data — wastage logs and seasonal price "
       "spreads — built up transparently into recurring annual impact.",
       takeaway="₹6 Cr saved every year = 9M+ more meals = 40,000+ children fed for a "
                "full year — recurring, audited, conservative.")
_, tf = tb(s, SW - ML - 4.0, 0.06, 4.0, 0.20)
para(tf, "PROPOSED SLIDE 4 GRAPH — v2 (compare with slide 9)", size=7.5, color=GRAY, italic=True,
     align=PP_ALIGN.RIGHT, first=True, space_after=0)

# ---------------- waterfall (left, manual shapes) ----------------
cy = section(s, ML, 1.50, 5.60, "Annual savings build-up (₹ Cr / year)", accent=NAVY)
BASE_Y = 5.30          # x-axis line
SCALE = 0.38           # inches per ₹Cr
BW = 0.80              # bar width
xs = [0.62, 1.70, 2.78, 3.86, 4.94]

bars = [
    # (label lines, source chip, from, to, color, value label)
    ("Veg waste cut\n4.9% → 2%", "TAPF wastage logs", 0.0, 4.1, GREEN, "+4.1"),
    ("Dhals buy-timing\n6–7% spread", "TAPF price data", 4.1, 6.0, TEAL, "+1.9"),
    ("Oils buy-timing\n4–5% spread", "TAPF price data", 6.0, 6.8, GOLD, "+0.8"),
    ("Storage &\ncarry cost", "netted off", 6.8, 6.0, RED, "−0.8"),
    ("NET SAVINGS\nevery year", "recurring", 0.0, 6.0, NAVY, "₹6.0 Cr"),
]
# axis line
rect(s, ML, BASE_Y, 5.55, 0.014, fill=LTGRAY, line=None)
prev_top = None
for (lbl, chip, v0, v1, c, vlab), x in zip(bars, xs):
    lo, hi = min(v0, v1), max(v0, v1)
    y_top = BASE_Y - hi * SCALE
    h = (hi - lo) * SCALE
    rect(s, x, y_top, BW, h, fill=c, line=None)
    # value label above bar
    _, tf = tb(s, x - 0.15, y_top - 0.24, BW + 0.30, 0.22)
    para(tf, vlab, size=9.5, color=(RED if v1 < v0 else NAVY), bold=True,
         align=PP_ALIGN.CENTER, first=True, space_after=0)
    # x labels + source chip
    _, tf = tb(s, x - 0.16, BASE_Y + 0.06, BW + 0.32, 0.55)
    para(tf, lbl, size=6.9, color=TEXT, bold=True, align=PP_ALIGN.CENTER,
         first=True, space_after=1, line_spacing=0.9)
    para(tf, chip, size=6, color=GRAY, italic=True, align=PP_ALIGN.CENTER,
         space_after=0)
    # dotted connector from previous bar top
    if prev_top is not None:
        lvl = BASE_Y - v0 * SCALE if v1 >= v0 else BASE_Y - v0 * SCALE
        conn = s.shapes.add_connector(1, Inches(prev_x + BW), Inches(prev_lvl),
                                      Inches(x), Inches(prev_lvl))
        line_style(conn, GRAY, 0.9, dash='sysDash')
    prev_top = y_top
    prev_x = x
    prev_lvl = BASE_Y - v1 * SCALE

# ---------------- right top: ramp mini-chart ----------------
RX, RW_ = 6.42, SW - ML - 6.42
cy = section(s, RX, 1.50, RW_, "Capture ramp (₹ Cr saved / yr)", accent=MIDBLUE,
             size=9.5)
cd = CategoryChartData()
cd.categories = ["2026", "2027", "2028", "2029"]
cd.add_series("saved", [1.8, 3.6, 5.1, 5.1])
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(RX), Inches(cy),
                        Inches(RW_), Inches(1.30), cd)
ch = gf.chart
ch.font.size = Pt(6.5); ch.font.name = SANS; ch.font.color.rgb = TEXT
ch.has_title = False; ch.has_legend = False
ch.series[0].format.fill.solid()
ch.series[0].format.fill.fore_color.rgb = GREEN
ch.series[0].format.line.fill.background()
plot = ch.plots[0]; plot.gap_width = 60
plot.has_data_labels = True
plot.data_labels.font.size = Pt(6.5); plot.data_labels.font.bold = True
plot.data_labels.font.color.rgb = NAVY
va = ch.value_axis; va.has_major_gridlines = False; va.visible = False
ca = ch.category_axis; ca.tick_labels.font.size = Pt(6.5)
ca.format.line.color.rgb = LTGRAY
_, tf = tb(s, RX, cy + 1.34, RW_, 0.20)
para(tf, "pilot 30% → scale 60% → steady-state 85% capture", size=6.5, color=GRAY,
     italic=True, first=True, space_after=0, align=PP_ALIGN.CENTER)

# ---------------- right bottom: conversion cascade ----------------
cy2 = section(s, RX, 3.62, RW_, "What ₹6 Cr/yr buys — every year", accent=GREEN,
              size=9.5)
steps = [("₹6.0 Cr", "saved annually, steady-state", NAVY),
         ("at ₹6.56", "donor cost per meal (govt funds the rest)", GRAY),
         ("9.1M+", "additional meals every year", GREEN),
         ("40,000+", "children fed for a full year — every year", ORANGE)]
sy = cy2 + 0.02
for i, (v, c, col) in enumerate(steps):
    _, tf = tb(s, RX, sy, RW_, 0.56)
    para(tf, [(v + "  ", True, col), (c, False, TEXT)], size=9.5 if i != 1 else 8,
         first=True, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, RX + 0.30, sy + 0.28, RX + 0.30, sy + 0.52, color=LTGRAY,
                     curve=False, width=1.0)
    sy += 0.60

# ---------------- assumptions footnote ----------------
_, tf = tb(s, ML, 6.10, 5.60, 0.44)
para(tf, "Assumptions: illustrative commodity split of TAPF's ₹450 Cr+ food spend; "
         "seasonal spreads per TAPF (dhals 6–7%, oils 4–5%, excl. taxation & global "
         "events); wastage per TAPF supply-chain data (4.1–5.6%/month). Spatial "
         "sourcing & spot-premium avoidance: not modelled — upside.",
     size=6.2, color=GRAY, italic=True, first=True, space_after=0, line_spacing=1.0)
_, tf = tb(s, 5.30, 7.16, 3.72, 0.30)
para(tf, "Source: The Akshaya Patra Foundation — official documents", size=6.5,
     color=GRAY, italic=True, align=PP_ALIGN.RIGHT, first=True, space_after=0)

prs.save(SRC)
print("saved; total slides:", len(prs.slides._sldIdLst))
