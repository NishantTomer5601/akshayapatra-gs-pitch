"""Append 4 alternative conclusion-slide layouts (slides 7-10) to the pitch deck."""
from deck_lib import *
from pptx import Presentation

PATH = "AkshayaPatra_GS_Pitch.pptx"
prs = Presentation(PATH)

TITLE = "Conclusion — why Akshaya Patra, and why this proposal wins"
SUB = ("A proven engine of social mobility, a grant that compounds instead of depletes, "
       "and a partnership only Goldman Sachs can power.")
TAKE = ("$250k that doesn't feed children for a year — it upgrades the engine "
        "that feeds them forever.")

def concl_slide(num, tag):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s, num, TITLE, SUB, takeaway=TAKE)
    _, tf = tb(s, SW - ML - 3.6, 0.06, 3.6, 0.20)
    para(tf, tag, size=7.5, color=GRAY, italic=True, align=PP_ALIGN.RIGHT, first=True,
         space_after=0)
    return s

WHY = [
    ("Structural, not sustenance",
     "others fund meals for a year; GS funds the technology that makes every future "
     "meal cheaper"),
    ("Plays to GS DNA",
     "forecasting, optimisation & risk controls — GS engineering mentors the build "
     "pro bono"),
    ("Measurable ROI",
     "$310k/yr audited savings; KPI dashboard shared quarterly with GS"),
    ("Scalable blueprint",
     "proven once, extends to school-feeding programmes across South Asia & Africa"),
]

QUOTE = ("“The mid-day meal was the reason I stayed in school. Today I manage "
         "portfolios instead of an empty stomach.”")
QUOTE_BY = [("— Krishna Kumar, ", True, NAVY),
            ("AVP, HSBC — former mid-day meal beneficiary", False, GRAY)]

RECOG = ["Padma Shri (2016) — Chairman Madhu Pandit Dasa",
         "World's largest NGO school-lunch programme — global record",
         "4-billionth meal celebrated with Hon'ble PM of India",
         "UK PM Rishi Sunak's kitchen visit — lauded hygiene & scale"]

PARTNERS = ["Infosys", "TCS", "HDFC", "Amazon", "Airbus", "SBI"]

# =====================================================================
# OPTION A — diagram-led, symmetric
# =====================================================================
s = concl_slide(7, "LAYOUT OPTION A — diagram-led, symmetric")

# centered equation hero
ov, oh = 1.32, 1.04
eq_w = 3*ov + 2*0.30 + 4*0.14
ex = (SW - eq_w) / 2
ey = 1.52
labels = [("GS $250k\ncatalytic\ncapital", PALEBLUE, NAVY),
          ("AP's 2.3M\nmeals-a-day\nengine", PALEGOLD, ORANGE),
          ("Impact that\ncompounds\nyearly", PALEGREEN, GREEN)]
cx = ex
for i, (t, bg, c) in enumerate(labels):
    rect(s, cx, ey, ov, oh, fill=bg, line=c, line_w=1.25, shape=MSO_SHAPE.OVAL)
    _, tf = tb(s, cx, ey + 0.14, ov, oh - 0.28, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7.8, color=c, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.92)
    if i < 2:
        sym = MSO_SHAPE.MATH_PLUS if i == 0 else MSO_SHAPE.MATH_EQUAL
        rect(s, cx + ov + 0.14, ey + oh/2 - 0.13, 0.30, 0.26, fill=NAVY, line=None,
             shape=sym)
    cx += ov + 0.30 + 0.28
_, tf = tb(s, 1.5, ey + oh + 0.08, SW - 3.0, 0.30)
para(tf, "One-time capital → permanent capability → recurring savings → more children "
         "in school, every year after Year 1", size=7.5, color=GRAY, first=True,
     space_after=0, align=PP_ALIGN.CENTER, italic=True)

# emphasized full-width WHY band
BY = 3.10; BH = 1.62
rect(s, ML, BY, SW - 2*ML, BH, fill=PALEBLUE, line=None)
rect(s, ML, BY, SW - 2*ML, 0.045, fill=NAVY, line=None)
_, tf = tb(s, ML, BY + 0.09, SW - 2*ML, 0.26)
para(tf, "WHY THE GOLDMAN SACHS GRANT STANDS OUT HERE", size=10.5, color=NAVY,
     bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
colw = (SW - 2*ML - 0.6) / 4
for i, (h, d) in enumerate(WHY):
    colx = ML + 0.30 + i * colw
    _, tf = tb(s, colx, BY + 0.42, colw - 0.18, BH - 0.52)
    para(tf, str(i + 1), size=15, color=BLUE, bold=True, first=True, space_after=1,
         font=SERIF)
    para(tf, h, size=8, color=NAVY, bold=True, space_after=1, line_spacing=0.95)
    para(tf, d, size=6.8, color=TEXT, space_after=0, line_spacing=0.95)

# bottom row: quote | recognition
vline(s, 4.98, 5.00, 6.45)
_, tf = tb(s, ML, 4.98, 4.35, 1.45)
para(tf, QUOTE, size=9, color=TEXT, italic=True, first=True, space_after=2,
     line_spacing=1.05)
para(tf, QUOTE_BY, size=7.5, space_after=4)
para(tf, [("92% secondary-school completion ", True, GREEN),
          ("among beneficiaries vs 74% national avg [dummy]", False, GRAY)],
     size=7.5, space_after=0)
_, tf = tb(s, 5.25, 4.98, SW - ML - 5.25, 1.45)
para(tf, "Recognised & trusted", size=8.5, color=NAVY, bold=True, first=True,
     space_after=2)
for r in RECOG[:3]:
    para(tf, "•  " + r, size=7.2, color=TEXT, space_after=1.5, line_spacing=0.98)
para(tf, [("Partners: ", True, NAVY),
          ("Infosys Foundation · TCS · HDFC · Amazon · Airbus · SBI + 200 more",
           False, GRAY)], size=7.2, space_after=0)

# =====================================================================
# OPTION B — typographic, ranked
# =====================================================================
s = concl_slide(8, "LAYOUT OPTION B — typographic, ranked")

# left big-numbers rail
railw = 2.60
rail = [("$250k", "one-time grant", NAVY, 24),
        ("$310k / yr", "returned in audited savings — every single year", GREEN, 19),
        ("+105,000", "children fed per year at zero extra cost", NAVY, 19),
        ("< 10 months", "payback on the full grant", ORANGE, 16)]
ry = 1.62
for i, (v, c, col, sz) in enumerate(rail):
    _, tf = tb(s, ML, ry, railw, 0.95)
    para(tf, v, size=sz, color=col, bold=True, font=SERIF, first=True, space_after=1)
    para(tf, c, size=7.5, color=GRAY, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, ML + 0.35, ry + 0.82, ML + 0.35, ry + 1.12, color=LTGRAY,
                     curve=False, width=1.0)
    ry += 1.22
vline(s, 3.15, 1.60, 6.45)

# right ranked rows
RXB = 3.45; RWB = SW - ML - RXB
ryy = 1.58
for i, (h, d) in enumerate(WHY):
    circ = rect(s, RXB, ryy, 0.34, 0.34, fill=NAVY if i < 2 else MIDBLUE, line=None,
                shape=MSO_SHAPE.OVAL)
    _, tf = tb(s, RXB, ryy + 0.015, 0.34, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, str(i + 1), size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    _, tf = tb(s, RXB + 0.48, ryy - 0.04, RWB - 0.48, 0.62)
    para(tf, h, size=9.5, color=NAVY, bold=True, first=True, space_after=1)
    para(tf, d, size=7.5, color=TEXT, space_after=0, line_spacing=0.98)
    ryy += 0.66

# comparison table
ty = ryy + 0.10
_, tf = tb(s, RXB, ty, RWB, 0.24)
para(tf, "A typical $250k grant  vs  this proposal", size=9, color=NAVY, bold=True,
     first=True, space_after=0)
rows = [("Feeds children for ~1 year", "Feeds 105,000 more children — every year"),
        ("Impact ends when funds end", "Savings outlive the cheque: 1.24x returned / yr"),
        ("Comes back for renewal next year", "Self-funding flywheel from month 10")]
tyy = ty + 0.30
for a, b in rows:
    rect(s, RXB, tyy + 0.30, RWB, 0.010, fill=RGBColor(0xE5, 0xE5, 0xE5), line=None)
    _, tf = tb(s, RXB, tyy + 0.02, RWB * 0.46, 0.28)
    para(tf, a, size=7.3, color=GRAY, first=True, space_after=0)
    _, tf = tb(s, RXB + RWB * 0.50, tyy + 0.02, RWB * 0.50, 0.28)
    para(tf, "✓  " + b, size=7.3, color=NAVY, bold=True, space_after=0, first=True)
    tyy += 0.36
_, tf = tb(s, RXB, tyy + 0.06, RWB, 0.60)
para(tf, [(QUOTE + "  ", False, TEXT)], size=7.3, italic=True, first=True,
     space_after=1, line_spacing=1.0)
para(tf, QUOTE_BY, size=7, space_after=0)

# =====================================================================
# OPTION C — timeline narrative
# =====================================================================
s = concl_slide(9, "LAYOUT OPTION C — timeline narrative")

# three columns above the spine
colx = [ML, 3.55, 6.70]
colw3 = 2.95
heads = [("PROVEN PAST", ORANGE), ("THE ASK — TODAY", NAVY),
         ("COMPOUNDING FUTURE", GREEN)]
# col 1: past
_, tf = tb(s, colx[0], 1.55, colw3, 2.05)
para(tf, "Two decades of delivery", size=9.5, color=NAVY, bold=True, first=True,
     space_after=3)
for r in RECOG[:3]:
    para(tf, "•  " + r, size=7.2, color=TEXT, space_after=1.5, line_spacing=0.98)
para(tf, [("“…I manage portfolios instead of an empty stomach.” ", False, TEXT),
          ("— Krishna Kumar, AVP HSBC, beneficiary", True, GRAY)], size=7,
     italic=True, space_after=0, line_spacing=0.98)
# col 2: ask card
card = rect(s, colx[1], 1.55, colw3, 2.02, fill=NAVY, line=None,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
_, tf = tb(s, colx[1] + 0.16, 1.68, colw3 - 0.32, 1.80)
para(tf, "$250,000", size=20, color=WHITE, bold=True, font=SERIF, first=True,
     space_after=2, align=PP_ALIGN.CENTER)
para(tf, "over 12 months · milestone-based tranches", size=7.5, color=WHITE,
     align=PP_ALIGN.CENTER, space_after=4)
para(tf, "•  Joint GS–AP steering committee", size=7.2, color=WHITE, space_after=1.5)
para(tf, "•  Quarterly KPI reviews with GS", size=7.2, color=WHITE, space_after=1.5)
para(tf, "•  Independent waste & delivery audits", size=7.2, color=WHITE,
     space_after=1.5)
para(tf, "•  GS volunteering & engineering mentorship", size=7.2, color=WHITE,
     space_after=0)
# col 3: future
_, tf = tb(s, colx[2], 1.55, colw3, 2.05)
para(tf, "What compounds after Year 1", size=9.5, color=NAVY, bold=True, first=True,
     space_after=3)
para(tf, [("$310k/yr ", True, GREEN), ("audited savings — every year, in perpetuity",
          False, TEXT)], size=7.4, space_after=1.5, line_spacing=0.98)
para(tf, [("+105,000 ", True, GREEN), ("children fed annually at zero extra donor cost",
          False, TEXT)], size=7.4, space_after=1.5, line_spacing=0.98)
para(tf, [("~70% ", True, GREEN), ("less food waste; 99%+ on-time delivery",
          False, TEXT)], size=7.4, space_after=1.5, line_spacing=0.98)
para(tf, [("Blueprint ", True, GREEN), ("extends to school-feeding programmes across "
          "South Asia & Africa", False, TEXT)], size=7.4, space_after=0,
     line_spacing=0.98)

# arrow spine
AY = 3.78
arr = rect(s, ML, AY, SW - 2*ML, 0.34, fill=LIGHTBLUE, line=None,
           shape=MSO_SHAPE.RIGHT_ARROW)
try:
    arr.adjustments[0] = 0.55; arr.adjustments[1] = 0.045
except Exception: pass
for (t, c), xx in zip(heads, colx):
    _, tf = tb(s, xx, AY + 0.055, colw3, 0.24)
    para(tf, t, size=8.5, color=c, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0)

# milestone track under spine
mstones = [("Mo 0", "grant signed"), ("Mo 4", "pilot live in 3 kitchens"),
           ("Mo 12", "both platforms in all 68 kitchens"),
           ("Yr 2", "flywheel self-funding"), ("Yr 3+", "blueprint scales abroad")]
mw = (SW - 2*ML) / 5
for i, (m, d) in enumerate(mstones):
    mx = ML + i * mw
    rect(s, mx + mw/2 - 0.008, 4.22, 0.016, 0.18, fill=GRAY, line=None)
    _, tf = tb(s, mx + 0.05, 4.44, mw - 0.10, 0.60)
    para(tf, m, size=8.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=1)
    para(tf, d, size=6.8, color=GRAY, align=PP_ALIGN.CENTER, space_after=0,
         line_spacing=0.92)

# bottom: why GS condensed + partners
rect(s, ML, 5.18, SW - 2*ML, 0.012, fill=RGBColor(0xE0, 0xE0, 0xE0), line=None)
_, tf = tb(s, ML, 5.30, SW - 2*ML, 0.55)
para(tf, [("Why GS stands out:  ", True, NAVY),
          ("structural fix, not sustenance  ·  plays to GS forecasting & optimisation "
           "DNA  ·  $310k/yr measurable ROI  ·  a blueprint the firm can scale globally",
           False, TEXT)], size=8, first=True, space_after=0, line_spacing=1.05,
     align=PP_ALIGN.CENTER)
_, tf = tb(s, ML, 5.92, SW - 2*ML, 0.30)
para(tf, [("In the company of: ", True, GRAY),
          ("Infosys Foundation · TCS · HDFC Bank · Amazon · Airbus · SBI Foundation "
           "+ 200 corporate donors", False, GRAY)], size=7.2, first=True,
     space_after=0, align=PP_ALIGN.CENTER, italic=True)

# =====================================================================
# OPTION D — sidebar emphasis
# =====================================================================
s = concl_slide(10, "LAYOUT OPTION D — sidebar emphasis")

# navy sidebar
SBW = 2.55
rect(s, ML, 1.50, SBW, 5.00, fill=NAVY, line=None)
_, tf = tb(s, ML + 0.18, 1.68, SBW - 0.36, 4.70)
para(tf, "THE ASK", size=9, color=GOLD, bold=True, first=True, space_after=2)
para(tf, "$250,000", size=23, color=WHITE, bold=True, font=SERIF, space_after=2)
para(tf, "over 12 months, in milestone-based tranches", size=7.5, color=WHITE,
     space_after=6, line_spacing=1.0)
para(tf, "Joint GS–AP steering committee  ·  quarterly KPI reviews  ·  independent "
         "waste & delivery audits", size=7.2, color=LIGHTBLUE, space_after=8,
     line_spacing=1.05)
para(tf, "WHAT IT RETURNS", size=9, color=GOLD, bold=True, space_after=2)
para(tf, "$310k / yr", size=16, color=WHITE, bold=True, space_after=1)
para(tf, "audited savings, every year", size=7.2, color=LIGHTBLUE, space_after=5)
para(tf, "+105,000", size=16, color=WHITE, bold=True, space_after=1)
para(tf, "children fed annually, zero extra cost", size=7.2, color=LIGHTBLUE,
     space_after=5)
para(tf, "< 10 months", size=16, color=WHITE, bold=True, space_after=1)
para(tf, "payback on the full grant", size=7.2, color=LIGHTBLUE, space_after=0)

# right content
RXD = ML + SBW + 0.45; RWD = SW - ML - RXD
_, tf = tb(s, RXD, 1.52, RWD, 0.30)
para(tf, "Why Goldman Sachs — and only Goldman Sachs", size=12, color=NAVY,
     bold=True, first=True, space_after=0)
rect(s, RXD, 1.84, RWD, 0.014, fill=LTGRAY, line=None)
yy = 1.96
for h, d in WHY:
    _, tf = tb(s, RXD, yy, RWD, 0.55)
    para(tf, [("✓  " + h + ":  ", True, GREEN), (d, False, TEXT)], size=8.3,
         first=True, space_after=0, line_spacing=1.0)
    yy += 0.52

# testimonial with photo
ty = yy + 0.12
img_placeholder(s, RXD, ty, 0.80, 0.95, "\U0001F4F7\nKrishna\nKumar")
_, tf = tb(s, RXD + 0.95, ty + 0.02, RWD - 0.95, 0.95)
para(tf, QUOTE, size=8, color=TEXT, italic=True, first=True, space_after=2,
     line_spacing=1.02)
para(tf, QUOTE_BY, size=7.2, space_after=0)

# recognition + partners strip
ry2 = ty + 1.10
rect(s, RXD, ry2, RWD, 0.012, fill=RGBColor(0xE0, 0xE0, 0xE0), line=None)
_, tf = tb(s, RXD, ry2 + 0.08, RWD, 0.50)
para(tf, [("Recognition:  ", True, NAVY),
          ("Padma Shri (2016) · world's largest school-lunch programme · 4-billionth "
           "meal with Hon'ble PM · UK PM Sunak's kitchen visit", False, TEXT)],
     size=7.2, first=True, space_after=0, line_spacing=1.02)
pw = (RWD - 5 * 0.06) / 6
for i, p in enumerate(PARTNERS):
    img_placeholder(s, RXD + i * (pw + 0.06), ry2 + 0.62, pw, 0.42, p)

prs.save(PATH)
print("saved 4 conclusion options; total slides:", len(prs.slides.__iter__.__self__._sldIdLst))
