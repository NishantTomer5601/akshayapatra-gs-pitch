"""Append 4 content-final slides (v2 of slides 1, 3, 4, 6) with verified TAPF data.
Existing slides are untouched. Sources footer on every new slide."""
from deck_lib import *
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import math

PATH = "AkshayaPatra_GS_Pitch.pptx"
prs = Presentation(PATH)

TEAM = ("Team: Nishant Tomer · Pratham Goenka · Aadithya Muralidharan · "
        "Somya Sethi · Mohajit Neog")

def new_slide(num_label, title, sub, takeaway, tag, team=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s, len(prs.slides._sldIdLst), title, sub, takeaway=takeaway)
    _, tf = tb(s, SW - ML - 4.4, 0.06, 4.4, 0.20)
    para(tf, tag, size=7.5, color=GRAY, italic=True, align=PP_ALIGN.RIGHT,
         first=True, space_after=0)
    _, tf = tb(s, 5.30, 7.16, 3.72, 0.30)
    para(tf, "Source: The Akshaya Patra Foundation — official documents",
         size=6.5, color=GRAY, italic=True, align=PP_ALIGN.RIGHT, first=True,
         space_after=0)
    if team:
        _, tf = tb(s, 1.72, 7.10, 3.45, 0.38)
        para(tf, TEAM, size=6.3, color=GRAY, first=True, space_after=0,
             line_spacing=0.95)
    return s

# =====================================================================
# SLIDE 7 — v2 of slide 1: introduction (verified numbers)
# =====================================================================
s = new_slide(7, "Who is The Akshaya Patra Foundation?",
              "The world's largest NGO-run school meal programme — serving 2.35 million "
              "children every school day, so that hunger never interrupts education.",
              "“No child in India shall be deprived of education because of hunger” — "
              "5 billion meals over 25 years.",
              "FINAL CONTENT — SLIDE 1 (v2)", team=True)

rect(s, ML, 1.42, SW - 2*ML, 0.80, fill=ORANGE, line=None)
img_placeholder(s, ML + 0.08, 1.51, 0.98, 0.62, "AKSHAYA PATRA\nlogo")
_, tf = tb(s, ML + 1.22, 1.42, SW - 2*ML - 1.36, 0.80, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Founded in 2000 serving 1,500 children, ", True, WHITE),
          ("Akshaya Patra is a Government-partnered PPP that has grown into the world's "
           "largest NGO-run school meal programme — ", False, WHITE),
          ("5 billion cumulative meals over 25 years", True, WHITE),
          (", a milestone celebrated with the Hon'ble President of India in March 2026.",
           False, WHITE)],
     size=8, first=True, space_after=0, line_spacing=1.0)

LX, LW = ML, 4.42
RX, RW = 5.02, SW - ML - 5.02
vline(s, 4.77, 2.42, 6.45)

cy = section(s, LX, 2.38, LW, "Who they are & why we chose them", accent=NAVY)
_, tf = tb(s, LX, cy, LW, 1.55)
para(tf, [("Mission: ", True, NAVY),
          ("serve mid-day meals to 3 million children every day and achieve 3 million "
           "servings of morning nutrition.", False, TEXT)], size=8, first=True,
     line_spacing=1.02)
para(tf, [("Why Akshaya Patra: ", True, NAVY),
          ("our supply-chain review found world-class production — and ", False, TEXT),
          ("two clearly identified, fixable gaps", True, TEXT),
          (" that TAPF itself ranks as its top technology priorities, where $250k "
           "unlocks compounding year-on-year impact.", False, TEXT)], size=8,
     line_spacing=1.02)
para(tf, "✓ 25 yrs of proven delivery   ✓ Govt-audited PPP   ✓ Tech-forward ops "
         "(ERP, OS1 dispatch, kitchen automation)",
     size=7.6, color=GREEN, bold=True, space_after=0)

cy = section(s, LX, 4.32, LW, "The journey — daily meals served (millions)", accent=BLUE)
line_chart(s, LX, cy + 0.02, LW, 1.42,
           ["2000", "2005", "2010", "2015", "2020", "2023", "2026"],
           [("Meals/day (M)", [0.0015, 0.33, 1.20, 1.44, 1.80, 2.00, 2.35])], [NAVY],
           font_size=7)
_, tf = tb(s, LX, cy + 1.50, LW, 0.34)
para(tf, [("Mar 2026: ", True, TEXT),
          ("5-billionth meal celebrated by the Hon'ble President of India  ·  ",
           False, TEXT),
          ("Apr 2026: ", True, TEXT),
          ("80th kitchen opens in Pune (Deutsche Bank CSR)", False, TEXT)],
     size=6.8, first=True, space_after=0, color=GRAY, line_spacing=0.95)

cy = section(s, RX, 2.38, RW, "What they do — the hub-and-spoke model", accent=MIDBLUE)
_, tf = tb(s, RX, cy, RW, 1.15)
para(tf, [("•  Centralised mega-kitchens ", True, TEXT),
          ("— rice for 1,000 students in 15 min; 40,000 rotis/hour", False, TEXT)],
     size=7.8, first=True, line_spacing=1.0)
para(tf, [("•  Hub-and-spoke delivery ", True, TEXT),
          ("— 1,100+ vehicles move 2 lakh+ insulated vessels daily", False, TEXT)],
     size=7.8, line_spacing=1.0)
para(tf, [("•  PPP economics ", True, TEXT),
          ("— govt funds ₹8.56 of every ₹17 meal; donors bridge ₹6.56", False, TEXT)],
     size=7.8, line_spacing=1.0)
para(tf, [("•  Beyond lunch ", True, TEXT),
          ("— morning nutrition, school rejuvenation, scholarships (NEST)",
           False, TEXT)], size=7.8, space_after=0, line_spacing=1.0)

cy = section(s, RX, 3.94, RW, "Scale at a glance (2026)", accent=ORANGE)
qw = RW / 4
for i, (v, c) in enumerate([("2.35M", "meals served\nevery day"),
                            ("80", "kitchens across\n81 locations"),
                            ("16+3", "states &\nunion territories"),
                            ("25,000+", "govt schools\nserved")]):
    stat(s, RX + i*qw, cy + 0.02, qw, v, c, vsize=14, csize=6.8)

cy = section(s, RX, 5.12, RW, "The opportunity — India context", accent=GREEN)
qw2 = 1.05
for i, (v, c) in enumerate([("2%", "of India's 130M govt-\nschool children reached"),
                            ("8.2%", "secondary dropout\nrate (UDISE+)"),
                            ("102/123", "India, Global Hunger\nIndex 2025")]):
    stat(s, RX + i * 1.18, cy + 0.02, 1.12, v, c, vsize=12, csize=6.2)
img_placeholder(s, RX + 3.60, cy + 0.04, RW - 3.60, 0.72,
                "\U0001F4F7 PHOTO: children\nat mid-day meal")

# =====================================================================
# SLIDE 8 — v2 of slide 3: gaps + $250k split (verified content)
# =====================================================================
s = new_slide(8, "The two gaps — and the $250,000 that closes them",
              "The gaps sit exactly where TAPF's own leadership says they are — its "
              "top-ranked technology priorities, confirmed in writing.",
              "Two fixable gaps, two targeted builds — both at the top of TAPF's own "
              "priority list.",
              "FINAL CONTENT — SLIDE 3 (v2)")

PANEL3 = RGBColor(0xF4, 0xF4, 0xF4)
QX3 = [ML, 5.08]; QW3 = 4.52
QY3 = [1.46, 4.02]; QH3 = [2.46, 2.44]
for gx in range(2):
    for gy in range(2):
        rect(s, QX3[gx], QY3[gy], QW3, QH3[gy], fill=PANEL3, line=None)

def qhead3(x, y, text, accent):
    rect(s, x, y + 0.03, 0.06, 0.20, fill=accent, line=None)
    _, tf = tb(s, x + 0.15, y, 3.2, 0.26)
    para(tf, text, size=10.5, color=NAVY, bold=True, first=True, space_after=0)

qhead3(QX3[0] + 0.15, QY3[0] + 0.12, "GAP 1 — Procurement", RED)
_, tf = tb(s, QX3[0] + 0.18, QY3[0] + 0.50, 3.05, 1.90)
para(tf, [("•  Periodic, manual planning", True, TEXT),
          (" vs demand that shifts daily — weather, off-season supply, attendance "
           "swings", False, TEXT)], size=7.8, first=True, line_spacing=1.02)
para(tf, [("•  No buffer day", True, TEXT),
          (" — a planning miss is a same-day miss for children", False, TEXT)],
     size=7.8, line_spacing=1.02)
para(tf, [("•  Complexity: ", True, TEXT),
          ("₹450 Cr+ food spend · 700 items · 655 vendors · 78 buyers, 145 approvers",
           False, TEXT)], size=7.8, line_spacing=1.02)
para(tf, [("Today: 4–5.6% monthly vegetable wastage · spot buys above seasonal "
           "price troughs", True, RED)], size=7.8, space_after=0, line_spacing=1.02)

qhead3(QX3[1] + 1.30, QY3[0] + 0.12, "GAP 2 — Distribution", RED)
_, tf = tb(s, QX3[1] + 1.32, QY3[0] + 0.50, 3.05, 1.90)
para(tf, [("•  No real-time vehicle location", True, TEXT),
          (" — schools phone kitchens to confirm deliveries", False, TEXT)],
     size=7.8, first=True, line_spacing=1.02)
para(tf, [("•  No maintenance or driver-safety data", True, TEXT),
          (" — issues surface only after incidents", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("•  4-hour cook-to-consume window", True, TEXT),
          (" at 65°C across ~70 km routes — zero slack", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("TAPF's #1 strategic headwind: last-mile traceability & food-safety risk",
           True, RED)], size=7.8, space_after=0, line_spacing=1.02)

qhead3(QX3[0] + 0.15, QY3[1] + 0.12, "SOLUTION 1 — “Annapurna AI” · $150k", NAVY)
_, tf = tb(s, QX3[0] + 0.18, QY3[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("AI demand forecasting & procurement planning across all 80 kitchens",
           False, TEXT)], size=7.8, first=True, line_spacing=1.02)
para(tf, [("How: ", True, NAVY),
          ("attendance, calendars, menus & weather → daily indents; buy-timing tuned "
           "to seasonal price troughs across 655 vendors", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("→  cuts veg wastage toward ~1–2%", True, GREEN)], size=8,
     line_spacing=1.02)
para(tf, [("→  savings recur every year (team est.)", True, GREEN)], size=8,
     space_after=0, line_spacing=1.02)

qhead3(QX3[1] + 1.30, QY3[1] + 0.12, "SOLUTION 2 — “Last-Mile Shield” · $100k", GREEN)
_, tf = tb(s, QX3[1] + 1.32, QY3[1] + 0.50, 3.05, 1.85)
para(tf, [("What: ", True, NAVY),
          ("AIS-140 GPS + BLE door sensors fleet-wide; geofenced delivery "
           "confirmation; smart locks on high-risk routes", False, TEXT)], size=7.8,
     first=True, line_spacing=1.02)
para(tf, [("How: ", True, NAVY),
          ("accelerates TAPF's own in-house blueprint (Traccar + DOUP) — auto "
           "delivery confirmation to schools", False, TEXT)], size=7.8,
     line_spacing=1.02)
para(tf, [("→  on-time within the 4-hr window", True, GREEN)], size=8,
     line_spacing=1.02)
para(tf, [("→  maintenance & driver-safety telemetry included", True, GREEN)],
     size=8, space_after=0, line_spacing=1.02)

ccx3, ccy3 = 5.0, 3.96
badge_d = 2.70
rect(s, ccx3 - badge_d/2, ccy3 - badge_d/2, badge_d, badge_d, fill=WHITE,
     line=LTGRAY, line_w=1.0, shape=MSO_SHAPE.OVAL)
cd3 = CategoryChartData()
cd3.categories = ["Solution 1 — $150k", "Solution 2 — $100k"]
cd3.add_series('s', [150, 100])
chart_w3 = 2.55
gf3 = s.shapes.add_chart(XL_CHART_TYPE.PIE,
                         Inches(ccx3 - chart_w3/2), Inches(ccy3 - chart_w3/2),
                         Inches(chart_w3), Inches(chart_w3), cd3)
ch3 = gf3.chart
ch3.font.size = Pt(9); ch3.font.name = SANS; ch3.has_title = False
ch3.has_legend = False
for i, c in enumerate([NAVY, GREEN]):
    pt = ch3.series[0].points[i]
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
    pt.format.line.color.rgb = WHITE; pt.format.line.width = Pt(2)
plot3 = ch3.plots[0]
plot3.has_data_labels = True
dl3 = plot3.data_labels
dl3.show_value = True
dl3.number_format = '"$"0"k"'; dl3.number_format_is_linked = False
dl3.font.size = Pt(11); dl3.font.bold = True; dl3.font.color.rgb = WHITE
cs3 = ch3._chartSpace
spPr3 = cs3.makeelement(qn('c:spPr'), {})
spPr3.append(spPr3.makeelement(qn('a:noFill'), {}))
ln3 = spPr3.makeelement(qn('a:ln'), {})
ln3.append(ln3.makeelement(qn('a:noFill'), {}))
spPr3.append(ln3)
ext3 = cs3.find(qn('c:externalData'))
if ext3 is not None:
    ext3.addprevious(spPr3)
else:
    cs3.append(spPr3)
pill_w, pill_h = 2.55, 0.42
rect(s, ccx3 - pill_w/2, ccy3 + badge_d/2 - 0.02, pill_w, pill_h, fill=WHITE,
     line=LTGRAY, line_w=0.9, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
_, tf = tb(s, ccx3 - pill_w/2, ccy3 + badge_d/2 + 0.005, pill_w, pill_h - 0.05,
           anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("THE $250k GRANT", True, NAVY)], size=8.5, first=True, space_after=0,
     align=PP_ALIGN.CENTER)
para(tf, "both gaps: TAPF's self-ranked top priorities", size=6.4, color=GRAY,
     align=PP_ALIGN.CENTER, space_after=0, italic=True)

# =====================================================================
# SLIDE 9 — v2 of slide 4: Solution 1 deep dive (verified content)
# =====================================================================
s = new_slide(9, "Solution 1 — the AI-enabled planning & procurement build ($150k)",
              "Connecting real-time demand signals to kitchen production and "
              "procurement — the technology that closes Gap 1.",
              "“Feeding more children with the same rupee” — AI plus GS capital "
              "makes it possible.",
              "FINAL CONTENT — SLIDE 4 (v2)")

P1_HEAD = "PILLAR 1 — Predictive Demand Forecasting"
P2_HEAD = "PILLAR 2 — Spatial-Temporal Procurement"
for x, headtxt, hc in [(ML, P1_HEAD, TEAL), (5.08, P2_HEAD, ORANGE)]:
    rect(s, x, 1.46, 4.52, 1.86, fill=WHITE, line=LTGRAY, line_w=0.9)
    rect(s, x, 1.46, 4.52, 0.30, fill=hc, line=None)
    _, tf = tb(s, x + 0.12, 1.49, 4.28, 0.25)
    para(tf, headtxt, size=9.5, color=WHITE, bold=True, first=True, space_after=0)
_, tf = tb(s, ML + 0.14, 1.86, 4.24, 1.42)
para(tf, "Predicts exact meal & ingredient needs per centre by ingesting:", size=7.4,
     color=GRAY, italic=True, first=True, space_after=2, line_spacing=1.0)
para(tf, [("•  Attendance & enrolment trends ", True, TEXT),
          ("estimate how many meals each centre should prepare daily", False, TEXT)],
     size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Holiday calendars ", True, TEXT),
          ("capture local demand spikes or dips before production is finalised",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Menu consumption history ", True, TEXT),
          ("converts expected meals into ingredient quantities for key staples",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Weather & consumption feedback ", True, TEXT),
          ("continuously refine the next centre-level forecast", False, TEXT)],
     size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Perishables (veg, dairy) ", True, TEXT),
          ("get shorter forecast cycles than staples (grains, dhals, oils)",
           False, TEXT)], size=7.4, space_after=0, line_spacing=1.02)
_, tf = tb(s, 5.22, 1.86, 4.24, 1.42)
para(tf, "Decides when and where to buy each commodity to minimise total landed cost:",
     size=7.4, color=GRAY, italic=True, first=True, space_after=2, line_spacing=1.0)
para(tf, [("•  Seasonal price troughs ", True, TEXT),
          ("identify when each commodity should be bought before prices rise",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Surplus-region signals ", True, TEXT),
          ("identify where each commodity can be sourced at lower landed cost",
           False, TEXT)], size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Logistics, handling & storage costs ", True, TEXT),
          ("are included in every sourcing recommendation", False, TEXT)],
     size=7.4, space_after=2, line_spacing=1.02)
para(tf, [("•  Quality grades & shelf-life rules ", True, TEXT),
          ("are locked into the final buying plan", False, TEXT)], size=7.4,
     space_after=2, line_spacing=1.02)
para(tf, [("•  Full procure-to-pay coverage ", True, TEXT),
          ("— requisition → PO → goods receipt → invoice → payment", False, TEXT)],
     size=7.4, space_after=0, line_spacing=1.02)

_, tf = tb(s, ML, 3.38, SW - 2*ML, 0.22)
para(tf, [("Annual volumes:  ", True, NAVY),
          ("120 T vegetables · 55 T dhals · 20 T dairy · 18 T oils · 6 T spices — "
           "against ₹450 Cr+ annual food spend", False, GRAY)], size=7.4, first=True,
     space_after=0, align=PP_ALIGN.CENTER)

cy = section(s, ML, 3.66, 4.52,
             "The self-learning operations loop — sharper every day", accent=NAVY,
             size=10)
LOOP4 = [("1 · INGEST", "pulls live prices,\nweather & attendance", NAVY),
         ("2 · FORECAST", "ingredient-level\ndemand per centre", MIDBLUE),
         ("3 · OPTIMIZE", "computes when &\nwhere to buy", GREEN),
         ("4 · RECOMMEND", "actionable buying\nsheet for kitchens", ORANGE),
         ("5 · LEARN", "outcome feedback\ntrains the model", MAGENTA)]
ccx4, ccy4, r4 = ML + 2.26, 5.32, 0.74
pts4 = []
for i, (t, sub, c) in enumerate(LOOP4):
    ang = math.radians(-90 + i * 72)
    nx = ccx4 + r4 * math.cos(ang); ny = ccy4 + r4 * math.sin(ang)
    rect(s, nx - 0.50, ny - 0.18, 1.00, 0.36, fill=c, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
    _, tf = tb(s, nx - 0.50, ny - 0.16, 1.00, 0.32, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    if i == 0:
        cxx, cyy = nx, ny - 0.40
    else:
        cxx, cyy = nx, ny + 0.40
    _, tf = tb(s, cxx - 0.65, cyy - 0.145, 1.30, 0.29, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, sub, size=6.2, color=GRAY, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    pts4.append((nx, ny))
for i in range(5):
    x1, y1 = pts4[i]; x2, y2 = pts4[(i + 1) % 5]
    mx, my = (x1 + x2)/2, (y1 + y2)/2
    vx, vy = mx - ccx4, my - ccy4
    nrm = math.hypot(vx, vy) or 1
    dotted_arrow(s, x1 + (x2 - x1)*0.34 + vx/nrm*0.09,
                 y1 + (y2 - y1)*0.34 + vy/nrm*0.09,
                 x1 + (x2 - x1)*0.66 + vx/nrm*0.09,
                 y1 + (y2 - y1)*0.66 + vy/nrm*0.09, color=GRAY, width=1.0)
_, tf = tb(s, ccx4 - 0.55, ccy4 - 0.10, 1.10, 0.22, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "↻  runs daily", size=6.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
     first=True, space_after=0)

cy = section(s, 5.08, 3.66, 4.52,
             "Rice procurement cost — with vs. without GS funding ($M)",
             accent=GREEN, size=9.5)
cd4 = CategoryChartData()
cd4.categories = ["2026", "2027", "2028", "2029"]
cd4.add_series("Without GS funding", [10.8, 11.4, 12.0, 12.6])
cd4.add_series("With GS funding", [9.4, 9.2, 9.1, 9.0])
gf4 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.08),
                         Inches(cy - 0.02), Inches(4.52), Inches(1.68), cd4)
ch4 = gf4.chart
ch4.font.size = Pt(7); ch4.font.name = SANS; ch4.font.color.rgb = TEXT
ch4.has_title = False
ch4.has_legend = True
ch4.legend.position = XL_LEGEND_POSITION.BOTTOM
ch4.legend.include_in_layout = False
ch4.legend.font.size = Pt(7)
for ser, c in zip(ch4.series, [ORANGE, GREEN]):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
    ser.format.line.fill.background()
plot4 = ch4.plots[0]
plot4.gap_width = 80; plot4.overlap = -10
va4 = ch4.value_axis
va4.has_major_gridlines = False
va4.tick_labels.font.size = Pt(6.5)
va4.tick_labels.number_format = '"$"0"M"'
va4.tick_labels.number_format_is_linked = False
ca4 = ch4.category_axis
ca4.tick_labels.font.size = Pt(7)
ca4.format.line.color.rgb = LTGRAY
rect(s, 5.08, cy + 1.74, 4.52, 0.66, fill=GREEN, line=None,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
_, tf = tb(s, 5.18, cy + 1.78, 4.32, 0.58, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "CUMULATIVE 2026–29 IMPACT", size=7, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, first=True, space_after=1)
para(tf, "~$10.1M saved  |  21.6% reduction", size=11.5, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, space_after=1)
para(tf, "team estimate on TAPF's ₹450 Cr+ annual food spend", size=6.6, color=WHITE,
     align=PP_ALIGN.CENTER, space_after=0, italic=True)

# =====================================================================
# SLIDE 10 — v2 of slide 6: conclusion (verified content)
# =====================================================================
s = new_slide(10, "Conclusion — why Akshaya Patra, and why this proposal wins",
              "A proven engine of social mobility, economics where every donor rupee "
              "is matched by government subsidy, and a grant that compounds.",
              "$250k that doesn't feed children for a year — it upgrades the engine "
              "that feeds them forever.",
              "FINAL CONTENT — SLIDE 6 (v2)")

RAILW = 2.55
_, tf = tb(s, ML, 1.52, RAILW, 0.24)
para(tf, "THE FUND & ITS IMPACT", size=9, color=NAVY, bold=True, first=True,
     space_after=0)
rect(s, ML, 1.80, RAILW, 0.014, fill=LTGRAY, line=None)
rail = [("$250k", "one-time grant ≈ ₹2.1 Cr", NAVY, 22),
        ("₹6.56", "donor cost per meal — govt funds ₹8.56 of every ₹17", GREEN, 20),
        ("₹1,500", "feeds one child for a full year (232 school days)", NAVY, 20),
        ("14,000+", "children fed for a year by the grant alone — before savings", ORANGE, 18)]
ry = 1.96
for i, (v, c, col, sz) in enumerate(rail):
    _, tf = tb(s, ML, ry, RAILW, 0.90)
    para(tf, v, size=sz, color=col, bold=True, font=SERIF, first=True, space_after=1)
    para(tf, c, size=7, color=GRAY, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, ML + 0.33, ry + 0.66, ML + 0.33, ry + 0.90, color=LTGRAY,
                     curve=False, width=1.0)
    ry += 0.95
vline(s, 3.05, 1.55, 5.32)

RX6, RW6 = 3.35, SW - ML - 3.35
cy = section(s, RX6, 1.50, RW6, "The partnership equation", accent=MIDBLUE, size=10.5)
ov, oh = 1.06, 0.88
eq_w = 3*ov + 2*0.50
ox1 = RX6 + (RW6 - eq_w) / 2
eq_y = cy + 0.02
for i, (t, bg, c) in enumerate([("GS $250k\ncatalytic\ncapital", PALEBLUE, NAVY),
                                ("AP's 2.35M\nmeals-a-day\nengine", PALEGOLD, ORANGE),
                                ("Impact that\ncompounds\nyearly", PALEGREEN, GREEN)]):
    xx = ox1 + i * (ov + 0.50)
    rect(s, xx, eq_y, ov, oh, fill=bg, line=c, line_w=1.1, shape=MSO_SHAPE.OVAL)
    _, tf = tb(s, xx, eq_y + 0.11, ov, oh - 0.22, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7, color=c, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    if i < 2:
        sym = MSO_SHAPE.MATH_PLUS if i == 0 else MSO_SHAPE.MATH_EQUAL
        rect(s, xx + ov + 0.12, eq_y + oh/2 - 0.12, 0.26, 0.24, fill=NAVY, line=None,
             shape=sym)
_, tf = tb(s, RX6, eq_y + oh + 0.05, RW6, 0.24)
para(tf, "One-time capital → permanent capability → recurring savings → more children "
         "in school, every year after Year 1", size=7, color=GRAY, italic=True,
     first=True, space_after=0, align=PP_ALIGN.CENTER)

cy2 = section(s, RX6, 3.16, RW6, "Why the GS grant stands out here", accent=NAVY,
              size=10.5)
WHY6 = [
    ("Structural, not sustenance", "TAPF's own words: restricted funds “deliver meals "
     "today” but cannot build “the systems behind every meal” — GS funds exactly those "
     "systems"),
    ("Co-designed priorities", "forecasting & real-time fleet visibility are TAPF's "
     "self-declared top tech gaps; logistics tech ranked #1 for grant ROI"),
    ("Plays to GS DNA", "forecasting, optimisation & risk controls — with GS "
     "engineering mentorship pro bono"),
    ("Measurable ROI", "audited savings against a ₹450 Cr+ procurement base; "
     "quarterly KPI dashboard to GS"),
    ("Compounding & green", "savings recur yearly and support TAPF's 50% "
     "renewable-by-2030 and EV-fleet-by-2035 goals"),
]
yy = cy2
for h, d in WHY6:
    _, tf = tb(s, RX6, yy, RW6, 0.42)
    para(tf, [("✓  " + h + ":  ", True, GREEN), (d, False, TEXT)], size=7.6,
         first=True, space_after=0, line_spacing=0.98)
    yy += 0.40

cy3 = section(s, ML, 5.44, SW - 2*ML, "Proven track record", accent=ORANGE, size=10.5)
TY = cy3 + 0.04
img_placeholder(s, ML, TY, 0.55, 0.62, "\U0001F4F7")
_, tf = tb(s, ML + 0.66, TY, 2.55, 0.85)
para(tf, "“Everybody used to get enough food and it was very tasty.”", size=6.9,
     color=TEXT, italic=True, first=True, space_after=1, line_spacing=1.0)
para(tf, [("— Krishna Kumar, ", True, NAVY),
          ("AVP, HSBC — programme alumnus", False, GRAY)], size=6.6, space_after=0)
vline(s, ML + 3.32, TY, TY + 0.72)
img_placeholder(s, ML + 3.48, TY, 0.55, 0.62, "\U0001F4F7")
_, tf = tb(s, ML + 4.14, TY, 2.45, 0.85)
para(tf, "“My dream is to become financially strong so I can support others the way "
         "I was supported.”", size=6.9, color=TEXT, italic=True, first=True,
     space_after=1, line_spacing=1.0)
para(tf, [("— Suma, ", True, NAVY),
          ("Purchasing Analyst, Palo Alto Networks", False, GRAY)], size=6.6,
     space_after=0)
vline(s, ML + 6.72, TY, TY + 0.72)
_, tf = tb(s, ML + 6.88, TY - 0.02, SW - ML - (ML + 6.88), 0.90)
para(tf, [("Recognition: ", True, NAVY),
          ("Padma Shri · Gandhi Peace Prize 2016 · BBC Global Champion · Nikkei Asia "
           "Award · CII Food Safety 2026 · HBS case study", False, TEXT)], size=6.6,
     first=True, space_after=2, line_spacing=1.0)
para(tf, [("Partners incl.: ", True, NAVY),
          ("Deutsche Bank (Pune kitchen '26) · IOCL (solar) + leading corporates",
           False, TEXT)], size=6.6, space_after=0, line_spacing=1.0)

prs.save(PATH)
print("saved; total slides:", len(prs.slides._sldIdLst))
