"""Append proposed slide-3 v2: gaps elaborated -> $250k split -> solutions introduced."""
from deck_lib import *
from pptx import Presentation

PATH = "AkshayaPatra_GS_Pitch.pptx"
prs = Presentation(PATH)

s = prs.slides.add_slide(prs.slide_layouts[6])
chrome(s, len(prs.slides._sldIdLst),
       "The two gaps — and the $250,000 that closes them",
       "Unpacking the gaps highlighted in the supply chain: what is broken today, how the "
       "grant splits across them, and the solution we build for each.",
       takeaway="Two fixable gaps, two targeted builds — $250k of capital becomes "
                "$310k of savings, every year.")
_, tf = tb(s, SW - ML - 3.6, 0.06, 3.6, 0.20)
para(tf, "PROPOSED SLIDE 3 — v2", size=7.5, color=GRAY, italic=True,
     align=PP_ALIGN.RIGHT, first=True, space_after=0)

# ---------------- allocation bar (top) ----------------
cy = section(s, ML, 1.45, SW - 2*ML, "How the $250,000 grant splits across the two gaps",
             accent=NAVY)
bar_y = cy + 0.06; bar_h = 0.46
bx = ML + 0.15; bw_total = SW - 2*ML - 0.3
w1 = bw_total * 0.6; w2 = bw_total * 0.4
rect(s, bx, bar_y, w1, bar_h, fill=NAVY, line=None)
rect(s, bx + w1, bar_y, w2, bar_h, fill=GREEN, line=None)
_, tf = tb(s, bx, bar_y, w1, bar_h, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "$150k — Gap 1: Procurement (60%)", size=10, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, first=True, space_after=0)
_, tf = tb(s, bx + w1, bar_y, w2, bar_h, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "$100k — Gap 2: Distribution (40%)", size=10, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, first=True, space_after=0)
_, tf = tb(s, bx, bar_y + bar_h + 0.04, w1, 0.24)
para(tf, "The costlier gap first — waste & over-priced buying start at the source",
     size=7.2, color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=0)
_, tf = tb(s, bx + w1, bar_y + bar_h + 0.04, w2, 0.24)
para(tf, "Then seal the last mile — on time, intact, untouched",
     size=7.2, color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=0)

# ---------------- two problem -> solution columns ----------------
C1X, C2X, CW = ML, 5.15, 4.45

# --- column 1: procurement ---
cy1 = section(s, C1X, 2.78, CW, "GAP 1 — Procurement: demand is guesswork",
              accent=RED, size=10.5)
_, tf = tb(s, C1X, cy1, CW, 1.10)
para(tf, [("•  Manual, gut-feel indenting", True, TEXT),
          (" 3 days ahead — blind to attendance swings, exams & holidays", False, TEXT)],
     size=8.2, first=True, line_spacing=1.05)
para(tf, [("•  Reactive spot-buying", True, TEXT),
          (" at mandi peak prices when stocks run short", False, TEXT)],
     size=8.2, line_spacing=1.05)
para(tf, [("•  No price intelligence", True, TEXT),
          (" — cheapest forward-buying windows are missed", False, TEXT)],
     size=8.2, line_spacing=1.05)
para(tf, [("•  No feedback loop", True, TEXT),
          (" from schools on actual meal uptake", False, TEXT)],
     size=8.2, space_after=0, line_spacing=1.05)
_, tf = tb(s, C1X, cy1 + 1.12, CW, 0.24)
para(tf, [("Cost today:  8–10% meals over-produced  ·  +12% paid on spot buys [dummy]",
           True, RED)], size=8.2, first=True, space_after=0)

dotted_arrow(s, C1X + CW/2, cy1 + 1.40, C1X + CW/2, cy1 + 1.60, color=NAVY,
             curve=False, width=1.25)

sb1y = cy1 + 1.66
rect(s, C1X, sb1y, CW, 1.52, fill=PALEBLUE, line=NAVY, line_w=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
_, tf = tb(s, C1X + 0.14, sb1y + 0.09, CW - 0.28, 1.52)
para(tf, [("SOLUTION 1 — “Annapurna AI”  ·  $150k", True, NAVY)], size=10.5,
     first=True, space_after=3)
para(tf, [("What: ", True, NAVY),
          ("ML demand-forecasting & procurement-planning platform across all 68 kitchens",
           False, TEXT)], size=8.2, line_spacing=1.05)
para(tf, [("How: ", True, NAVY),
          ("predicts school-level attendance from history, calendars & seasonality; "
           "auto-generates daily indents & forward buying schedules", False, TEXT)],
     size=8.2, line_spacing=1.05)
para(tf, [("Rollout: ", True, GRAY),
          ("3-kitchen pilot (Mo 1–4) → all 68 kitchens by Mo 12", False, GRAY)],
     size=7.8, line_spacing=1.05)
para(tf, [("→  ~70% less food waste  ·  est. $230k saved / yr [dummy]", True, GREEN)],
     size=8.6, space_after=0)

# --- column 2: distribution ---
cy2 = section(s, C2X, 2.78, CW, "GAP 2 — Distribution: the last mile is blind",
              accent=RED, size=10.5)
_, tf = tb(s, C2X, cy2, CW, 1.10)
para(tf, [("•  Static, driver-memory routes", True, TEXT),
          (" — 15–30 min past the lunch bell means children stay hungry all day",
           False, TEXT)], size=8.2, first=True, line_spacing=1.05)
para(tf, [("•  Untracked vessels & utensils", True, TEXT),
          (" — daily losses add to operating cost", False, TEXT)],
     size=8.2, line_spacing=1.05)
para(tf, [("•  Adulteration / diversion risk", True, TEXT),
          (" — unsealed vans can be offloaded at roadside dhabas", False, TEXT)],
     size=8.2, line_spacing=1.05)
para(tf, [("•  No proof-of-delivery", True, TEXT),
          (" — school shortfall disputes go unresolved", False, TEXT)],
     size=8.2, space_after=0, line_spacing=1.05)
_, tf = tb(s, C2X, cy2 + 1.12, CW, 0.24)
para(tf, [("Cost today:  6.2% deliveries miss lunch  ·  $35k/yr shrinkage [dummy]",
           True, RED)], size=8.2, first=True, space_after=0)

dotted_arrow(s, C2X + CW/2, cy2 + 1.40, C2X + CW/2, cy2 + 1.60, color=GREEN,
             curve=False, width=1.25)

sb2y = cy2 + 1.66
rect(s, C2X, sb2y, CW, 1.52, fill=PALEGREEN, line=GREEN, line_w=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
_, tf = tb(s, C2X + 0.14, sb2y + 0.09, CW - 0.28, 1.52)
para(tf, [("SOLUTION 2 — “Last-Mile Shield”  ·  $100k", True, GREEN)], size=10.5,
     first=True, space_after=3)
para(tf, [("What: ", True, NAVY),
          ("AI route optimisation + van telematics (GPS, cameras, RFID tags) + "
           "geofenced smart locks", False, TEXT)], size=8.2, line_spacing=1.05)
para(tf, [("How: ", True, NAVY),
          ("routes sequenced to each school's lunch bell; doors unlock only within "
           "150 m of a registered school — diversion becomes impossible", False, TEXT)],
     size=8.2, line_spacing=1.05)
para(tf, [("Bonus: ", True, GRAY),
          ("hardware amortises over 5+ yrs; 3 vans freed per hub reach new schools",
           False, GRAY)], size=7.8, line_spacing=1.05)
para(tf, [("→  99%+ on-time  ·  −18% fleet km  ·  est. $80k saved / yr [dummy]",
           True, GREEN)], size=8.6, space_after=0)

prs.save(PATH)
print("saved; total slides:", len(prs.slides._sldIdLst))
