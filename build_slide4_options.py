"""Append 4 layout options for slide 4 (Solution 1 deep-dive) as new slides.
Content: two pillars, self-learning ops loop, rice cost chart, $10.1M impact."""
from deck_lib import *
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import math

PATH = "AkshayaPatra_GS_Pitch.pptx"
prs = Presentation(PATH)

TITLE = "Solution 1 — the AI-enabled planning & procurement build ($150k)"
SUB = ("Connecting real-time demand signals to kitchen production and procurement — "
       "the technology that closes Gap 1.")
TAKE = ("“Feeding more children with the same rupee” — AI plus GS capital makes it "
        "possible.")

P1_HEAD = "PILLAR 1 — Predictive Demand Forecasting"
P1_INTRO = "Predicts exact meal & ingredient needs per centre by ingesting:"
P1 = [
    ("Attendance & enrolment trends", "estimate how many meals each centre should "
     "prepare daily"),
    ("Holiday calendars", "capture local demand spikes or dips before production "
     "quantities are finalised"),
    ("Menu consumption history", "converts expected meals into ingredient quantities "
     "for key staples"),
    ("Weather signals & feedback", "continuously refine the next centre-level forecast"),
]
P2_HEAD = "PILLAR 2 — Spatial-Temporal Procurement"
P2_INTRO = "Decides when and where to buy each commodity to minimise total landed cost:"
P2 = [
    ("Seasonal price troughs", "identify when each commodity should be bought before "
     "prices rise"),
    ("Surplus-region signals", "identify where each commodity can be sourced at lower "
     "landed cost"),
    ("Logistics, handling & storage costs", "are included in every sourcing "
     "recommendation"),
    ("Quality grades & shelf-life rules", "are locked into the final buying plan"),
]
LOOP = [("INGEST", "live prices, weather,\nattendance", NAVY),
        ("FORECAST", "ingredient-level\ndemand", MIDBLUE),
        ("OPTIMIZE", "when & where\nto buy", GREEN),
        ("RECOMMEND", "actionable\nbuying sheet", ORANGE),
        ("LEARN", "feedback trains\nthe model", MAGENTA)]
YEARS = ["2026", "2027", "2028", "2029"]
WITHOUT = [10.8, 11.4, 12.0, 12.6]
WITH = [9.4, 9.2, 9.1, 9.0]
CHART_HEAD = "Rice procurement cost — with vs. without GS funding ($M)"

def new_option(tag):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s, len(prs.slides._sldIdLst), TITLE, SUB, takeaway=TAKE)
    _, tf = tb(s, SW - ML - 4.4, 0.06, 4.4, 0.20)
    para(tf, tag, size=7.5, color=GRAY, italic=True, align=PP_ALIGN.RIGHT,
         first=True, space_after=0)
    return s

def pillar_bullets(s, x, y, w, items, intro=None, size=7.8, gap=2.5):
    _, tf = tb(s, x, y, w, 2.2)
    if intro:
        para(tf, intro, size=size, color=GRAY, italic=True, first=True,
             space_after=gap, line_spacing=1.0)
    for i, (h, d) in enumerate(items):
        para(tf, [("•  " + h + " ", True, TEXT), (d, False, TEXT)], size=size,
             first=(i == 0 and not intro), space_after=gap, line_spacing=1.02)
    return tf

def grouped_chart(s, x, y, w, h, labels=False, legend=True, font=7):
    cd = CategoryChartData()
    cd.categories = YEARS
    cd.add_series("Without GS funding", WITHOUT)
    cd.add_series("With GS funding", WITH)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y),
                            Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.font.size = Pt(font); ch.font.name = SANS; ch.font.color.rgb = TEXT
    ch.has_title = False
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(font)
    for ser, c in zip(ch.series, [ORANGE, GREEN]):
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
        ser.format.line.fill.background()
    plot = ch.plots[0]
    plot.gap_width = 80; plot.overlap = -10
    if labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(font - 0.5); dl.font.bold = True; dl.font.color.rgb = NAVY
        dl.number_format = '"$"0.0"M"'; dl.number_format_is_linked = False
    va = ch.value_axis
    va.has_major_gridlines = False; va.visible = True
    va.tick_labels.font.size = Pt(font - 0.5)
    va.tick_labels.number_format = '"$"0"M"'; va.tick_labels.number_format_is_linked = False
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(font)
    ca.format.line.color.rgb = LTGRAY
    return ch

def impact_band(s, x, y, w, h, big=12.5):
    rect(s, x, y, w, h, fill=GREEN, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10)
    _, tf = tb(s, x + 0.10, y + 0.04, w - 0.20, h - 0.08, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "CUMULATIVE 2026–29 IMPACT", size=7, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, first=True, space_after=1)
    para(tf, "~$10.1M saved  |  21.6% reduction", size=big, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, space_after=1)
    para(tf, "freed capital feeds more children", size=7, color=WHITE,
         align=PP_ALIGN.CENTER, space_after=0, italic=True)

def loop_diagram(s, ccx, ccy, r, node_w, node_h, tsize=7.2, ssize=None,
                 center_label=True):
    pts = []
    for i, (t, sub, c) in enumerate(LOOP):
        ang = math.radians(-90 + i * 72)
        nx = ccx + r * math.cos(ang); ny = ccy + r * math.sin(ang)
        rect(s, nx - node_w/2, ny - node_h/2, node_w, node_h, fill=c, line=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
        _, tf = tb(s, nx - node_w/2, ny - node_h/2 + 0.015, node_w, node_h - 0.03,
                   anchor=MSO_ANCHOR.MIDDLE)
        para(tf, t, size=tsize, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=0, line_spacing=0.9)
        if ssize:
            para(tf, sub, size=ssize, color=WHITE, align=PP_ALIGN.CENTER,
                 space_after=0, line_spacing=0.85)
        pts.append((nx, ny))
    for i in range(5):
        x1, y1 = pts[i]; x2, y2 = pts[(i + 1) % 5]
        mx, my = (x1 + x2)/2, (y1 + y2)/2
        vx, vy = mx - ccx, my - ccy
        nrm = math.hypot(vx, vy) or 1
        fx1 = x1 + (x2 - x1)*0.32 + vx/nrm*0.10
        fy1 = y1 + (y2 - y1)*0.32 + vy/nrm*0.10
        fx2 = x1 + (x2 - x1)*0.68 + vx/nrm*0.10
        fy2 = y1 + (y2 - y1)*0.68 + vy/nrm*0.10
        dotted_arrow(s, fx1, fy1, fx2, fy2, color=GRAY, width=1.0)
    if center_label:
        _, tf = tb(s, ccx - 0.75, ccy - 0.22, 1.5, 0.44, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "SELF-LEARNING\nOPS LOOP", size=7.5, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, first=True, space_after=0, line_spacing=0.95)

# =====================================================================
# OPTION A — faithful quadrant: pillars top, loop + chart bottom
# =====================================================================
s = new_option("PROPOSED SLIDE 4 — OPTION A (quadrant)")
for x, head, hc in [(ML, P1_HEAD, TEAL), (5.08, P2_HEAD, ORANGE)]:
    rect(s, x, 1.46, 4.52, 1.66, fill=WHITE, line=LTGRAY, line_w=0.9)
    rect(s, x, 1.46, 4.52, 0.30, fill=hc, line=None)
    _, tf = tb(s, x + 0.12, 1.49, 4.28, 0.25)
    para(tf, head, size=9.5, color=WHITE, bold=True, first=True, space_after=0)
pillar_bullets(s, ML + 0.14, 1.86, 4.24, P1, intro=P1_INTRO, size=7.5, gap=2)
pillar_bullets(s, 5.22, 1.86, 4.24, P2, intro=P2_INTRO, size=7.5, gap=2)

cy = section(s, ML, 3.30, 4.52, "The self-learning operations loop", accent=NAVY,
             size=10.5)
loop_diagram(s, ML + 2.26, 5.12, 0.95, 1.04, 0.42, tsize=7.2)

cy = section(s, 5.08, 3.30, 4.52, CHART_HEAD, accent=GREEN, size=9.5)
grouped_chart(s, 5.08, cy - 0.02, 4.52, 1.95, labels=False, legend=True, font=7)
impact_band(s, 5.08, cy + 2.00, 4.52, 0.68, big=11.5)

# =====================================================================
# OPTION B — pipeline spine: pillars top, loop as horizontal flow, results bottom
# =====================================================================
s = new_option("PROPOSED SLIDE 4 — OPTION B (pipeline)")
cy1 = section(s, ML, 1.45, 4.5, P1_HEAD, accent=TEAL, size=10)
pillar_bullets(s, ML, cy1, 4.5, P1, intro=None, size=7.5, gap=1.5)
cy2 = section(s, 5.1, 1.45, 4.5, P2_HEAD, accent=ORANGE, size=10)
pillar_bullets(s, 5.1, cy2, 4.5, P2, intro=None, size=7.5, gap=1.5)

# pipeline
PIPE_Y = 3.02; NW, NH, NGAP = 1.62, 0.62, 0.275
px = ML
for i, (t, sub, c) in enumerate(LOOP):
    rect(s, px, PIPE_Y, NW, NH, fill=c, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    _, tf = tb(s, px, PIPE_Y + 0.04, NW, NH - 0.08, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=1, line_spacing=0.9)
    para(tf, sub.replace("\n", " "), size=6, color=WHITE, align=PP_ALIGN.CENTER,
         space_after=0, line_spacing=0.9)
    if i < 4:
        dotted_arrow(s, px + NW + 0.03, PIPE_Y + NH/2, px + NW + NGAP - 0.03,
                     PIPE_Y + NH/2, color=GRAY, width=1.1)
    px += NW + NGAP
dotted_arrow(s, ML + 4*(NW+NGAP) + NW/2, PIPE_Y + NH + 0.05,
             ML + NW/2, PIPE_Y + NH + 0.22, color=GRAY, width=1.1)
_, tf = tb(s, ML, PIPE_Y + NH + 0.24, SW - 2*ML, 0.20)
para(tf, "…the loop repeats every day — each cycle makes the next forecast sharper",
     size=7, color=GRAY, italic=True, first=True, space_after=0, align=PP_ALIGN.CENTER)

cyc = section(s, ML, 4.40, 5.5, CHART_HEAD, accent=GREEN, size=9.5)
grouped_chart(s, ML, cyc - 0.02, 5.5, 1.62, labels=True, legend=True, font=6.8)
rect(s, 6.15, 4.40, 3.45, 2.02, fill=GREEN, line=None,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
_, tf = tb(s, 6.30, 4.56, 3.15, 1.80)
para(tf, "CUMULATIVE 2026–29 IMPACT", size=7.5, color=WHITE, bold=True, first=True,
     space_after=4)
para(tf, "~$10.1M", size=20, color=WHITE, bold=True, font=SERIF, space_after=1)
para(tf, "saved on procurement", size=7.5, color=WHITE, space_after=5)
para(tf, "21.6%", size=16, color=WHITE, bold=True, font=SERIF, space_after=1)
para(tf, "reduction in cost — freed capital feeds more children", size=7.5,
     color=WHITE, space_after=0, line_spacing=1.0)

# =====================================================================
# OPTION C — numbers rail: impact left, pillars + loop chips + chart right
# =====================================================================
s = new_option("PROPOSED SLIDE 4 — OPTION C (numbers rail)")
RAILW = 2.55
_, tf = tb(s, ML, 1.52, RAILW, 0.24)
para(tf, "THE RETURN ON $150k", size=9, color=NAVY, bold=True, first=True,
     space_after=0)
rect(s, ML, 1.80, RAILW, 0.014, fill=LTGRAY, line=None)
rail = [("$150k", "one-time GS-funded build", NAVY, 22),
        ("~$10.1M", "saved on procurement, cumulative 2026–29", GREEN, 20),
        ("21.6%", "reduction in commodity cost", GREEN, 20),
        ("100%", "of savings reinvested in more meals", ORANGE, 18)]
ry = 1.98
for i, (v, c, col, sz) in enumerate(rail):
    _, tf = tb(s, ML, ry, RAILW, 0.95)
    para(tf, v, size=sz, color=col, bold=True, font=SERIF, first=True, space_after=1)
    para(tf, c, size=7.4, color=GRAY, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, ML + 0.33, ry + 0.78, ML + 0.33, ry + 1.06, color=LTGRAY,
                     curve=False, width=1.0)
    ry += 1.16
vline(s, 3.05, 1.55, 6.45)

RX, RW = 3.35, SW - ML - 3.35
cyp = section(s, RX, 1.45, RW, P1_HEAD, accent=TEAL, size=10)
pillar_bullets(s, RX, cyp, RW, P1, intro=None, size=7.5, gap=1.5)
cyp = section(s, RX, 2.88, RW, P2_HEAD, accent=ORANGE, size=10)
pillar_bullets(s, RX, cyp, RW, P2, intro=None, size=7.5, gap=1.5)

# loop chips
chip_y = 4.42; cw_ = 1.10; cgap = 0.135
cx0 = RX
for i, (t, sub, c) in enumerate(LOOP):
    rect(s, cx0, chip_y, cw_, 0.34, fill=c, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    _, tf = tb(s, cx0, chip_y + 0.015, cw_, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0)
    if i < 4:
        dotted_arrow(s, cx0 + cw_ + 0.02, chip_y + 0.17, cx0 + cw_ + cgap - 0.02,
                     chip_y + 0.17, color=GRAY, width=1.0)
    cx0 += cw_ + cgap
_, tf = tb(s, RX, chip_y + 0.38, RW, 0.18)
para(tf, "the self-learning loop — repeats daily, each cycle sharpens the next "
         "forecast", size=6.8, color=GRAY, italic=True, first=True, space_after=0,
     align=PP_ALIGN.CENTER)

cyc = section(s, RX, 5.02, RW, CHART_HEAD, accent=GREEN, size=9.5)
grouped_chart(s, RX, cyc - 0.04, 4.15, 1.44, labels=False, legend=True, font=6.5)
_, tf = tb(s, RX + 4.30, cyc + 0.10, RW - 4.30, 1.30)
para(tf, "~$10.1M saved", size=13, color=GREEN, bold=True, font=SERIF, first=True,
     space_after=1)
para(tf, "21.6% ↓ cost, 2026–29", size=9, color=GREEN, bold=True, space_after=3)
para(tf, "freed capital feeds more children", size=7.2, color=GRAY, italic=True,
     space_after=0, line_spacing=1.0)

# =====================================================================
# OPTION D — central hub: big loop centre, pillars flanking, results bottom
# =====================================================================
s = new_option("PROPOSED SLIDE 4 — OPTION D (central hub)")
cyp = section(s, ML, 1.50, 2.95, P1_HEAD, accent=TEAL, size=9)
pillar_bullets(s, ML, cyp, 2.95, P1, intro=None, size=7.2, gap=2)
cyp = section(s, 6.65, 1.50, 2.95, P2_HEAD, accent=ORANGE, size=9)
pillar_bullets(s, 6.65, cyp, 2.95, P2, intro=None, size=7.2, gap=2)

ccx, ccy, r = 5.0, 3.10, 1.02
oval_d = 1.30
rect(s, ccx - oval_d/2, ccy - oval_d/2, oval_d, oval_d, fill=PALEBLUE, line=NAVY,
     line_w=1.1, shape=MSO_SHAPE.OVAL)
loop_diagram(s, ccx, ccy, r, 1.06, 0.56, tsize=7, ssize=5.5, center_label=False)
_, tf = tb(s, ccx - 0.62, ccy - 0.20, 1.24, 0.40, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "SELF-LEARNING\nOPS LOOP", size=6.8, color=NAVY, bold=True,
     align=PP_ALIGN.CENTER, first=True, space_after=0, line_spacing=0.95)
dotted_arrow(s, ML + 2.95 + 0.06, 2.55, ccx - r - 0.60, 2.95, color=TEAL, width=1.1)
dotted_arrow(s, 6.65 - 0.06, 2.55, ccx + r + 0.60, 2.95, color=ORANGE, width=1.1)

cyc = section(s, ML, 4.78, 5.45, CHART_HEAD, accent=GREEN, size=9.5)
grouped_chart(s, ML, cyc - 0.04, 5.45, 1.36, labels=False, legend=True, font=6.5)
impact_band(s, 6.10, 4.78, 3.50, 1.74, big=14)

prs.save(PATH)
print("saved; total slides:", len(prs.slides._sldIdLst))
