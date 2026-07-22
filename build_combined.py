"""Append Option E (slide 11): combined conclusion — fund-impact rail + equation +
why-GS + testimonial + corporate trust."""
from deck_lib import *
from pptx import Presentation

PATH = "AkshayaPatra_GS_Pitch.pptx"
prs = Presentation(PATH)

TITLE = "Conclusion — why Akshaya Patra, and why this proposal wins"
SUB = ("A proven engine of social mobility, a grant that compounds instead of depletes, "
       "and a partnership only Goldman Sachs can power.")
TAKE = ("$250k that doesn't feed children for a year — it upgrades the engine "
        "that feeds them forever.")

s = prs.slides.add_slide(prs.slide_layouts[6])
chrome(s, 11, TITLE, SUB, takeaway=TAKE)
_, tf = tb(s, SW - ML - 3.6, 0.06, 3.6, 0.20)
para(tf, "LAYOUT OPTION E — combined", size=7.5, color=GRAY, italic=True,
     align=PP_ALIGN.RIGHT, first=True, space_after=0)

# ---------------- left: fund & impact rail ----------------
RAILW = 2.55
_, tf = tb(s, ML, 1.52, RAILW, 0.24)
para(tf, "THE FUND & ITS IMPACT", size=9, color=NAVY, bold=True, first=True,
     space_after=0)
rect(s, ML, 1.80, RAILW, 0.014, fill=LTGRAY, line=None)

rail = [("$250k", "one-time grant", NAVY, 23),
        ("$310k / yr", "returned in audited savings — every single year", GREEN, 18),
        ("+105,000", "children fed per year at zero extra cost", NAVY, 18),
        ("< 10 months", "payback on the full grant", ORANGE, 15)]
ry = 1.98
for i, (v, c, col, sz) in enumerate(rail):
    _, tf = tb(s, ML, ry, RAILW, 0.92)
    para(tf, v, size=sz, color=col, bold=True, font=SERIF, first=True, space_after=1)
    para(tf, c, size=7.4, color=GRAY, space_after=0, line_spacing=0.95)
    if i < 3:
        dotted_arrow(s, ML + 0.33, ry + 0.76, ML + 0.33, ry + 1.04, color=LTGRAY,
                     curve=False, width=1.0)
    ry += 1.14

vline(s, 3.05, 1.55, 6.45)

# ---------------- right top: partnership equation ----------------
RX, RW = 3.35, SW - ML - 3.35
cy = section(s, RX, 1.50, RW, "The partnership equation", accent=MIDBLUE, size=10.5)
ov, oh = 1.06, 0.90
eq_w = 3*ov + 2*0.50
ox1 = RX + (RW - eq_w) / 2
eq_y = cy + 0.02
eqdata = [("GS $250k\ncatalytic\ncapital", PALEBLUE, NAVY),
          ("AP's 2.3M\nmeals-a-day\nengine", PALEGOLD, ORANGE),
          ("Impact that\ncompounds\nyearly", PALEGREEN, GREEN)]
xx = ox1
for i, (t, bg, c) in enumerate(eqdata):
    rect(s, xx, eq_y, ov, oh, fill=bg, line=c, line_w=1.1, shape=MSO_SHAPE.OVAL)
    _, tf = tb(s, xx, eq_y + 0.11, ov, oh - 0.22, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=7, color=c, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=0.9)
    if i < 2:
        sym = MSO_SHAPE.MATH_PLUS if i == 0 else MSO_SHAPE.MATH_EQUAL
        rect(s, xx + ov + 0.12, eq_y + oh/2 - 0.12, 0.26, 0.24, fill=NAVY, line=None,
             shape=sym)
    xx += ov + 0.50
_, tf = tb(s, RX, eq_y + oh + 0.06, RW, 0.26)
para(tf, "One-time capital → permanent capability → recurring savings → more children "
         "in school, every year after Year 1", size=7, color=GRAY, italic=True,
     first=True, space_after=0, align=PP_ALIGN.CENTER)

# ---------------- right middle: why GS ----------------
cy2 = section(s, RX, 3.32, RW, "Why the GS grant stands out here", accent=NAVY,
              size=10.5)
WHY = [
    ("Structural, not sustenance", "others fund meals for a year; GS funds the "
     "technology that makes every future meal cheaper"),
    ("Plays to GS DNA", "forecasting, optimisation & risk controls — GS engineering "
     "mentors the build pro bono"),
    ("Measurable ROI", "$310k/yr audited savings; KPI dashboard shared quarterly "
     "with GS"),
    ("Scalable blueprint", "proven once, extends to school-feeding programmes across "
     "South Asia & Africa"),
]
yy = cy2
for h, d in WHY:
    _, tf = tb(s, RX, yy, RW, 0.42)
    para(tf, [("✓  " + h + ":  ", True, GREEN), (d, False, TEXT)], size=8,
         first=True, space_after=0, line_spacing=0.98)
    yy += 0.385

# ---------------- right bottom: testimonial | trusted by ----------------
BY = yy + 0.14
rect(s, RX, BY, RW, 0.012, fill=RGBColor(0xE0, 0xE0, 0xE0), line=None)
TY = BY + 0.10
# testimonial (left half)
img_placeholder(s, RX, TY, 0.72, 0.92, "\U0001F4F7\nKrishna\nKumar")
_, tf = tb(s, RX + 0.86, TY + 0.02, 2.42, 0.95)
para(tf, "“The mid-day meal was the reason I stayed in school. Today I manage "
         "portfolios instead of an empty stomach.”", size=7.2, color=TEXT, italic=True,
     first=True, space_after=2, line_spacing=1.0)
para(tf, [("— Krishna Kumar, ", True, NAVY),
          ("AVP, HSBC — former beneficiary", False, GRAY)], size=6.8, space_after=0)
# divider
vline(s, RX + 3.46, TY, TY + 0.95)
# trusted by (right half)
PX = RX + 3.62; PW = RW - 3.62
_, tf = tb(s, PX, TY, PW, 0.22)
para(tf, [("Trusted by leading corporates ", True, NAVY), ("(+200 more)", False, GRAY)],
     size=8.2, first=True, space_after=0)
lw = (PW - 2 * 0.06) / 3
for i, p in enumerate(["Infosys", "TCS", "HDFC", "Amazon", "Airbus", "SBI"]):
    px = PX + (i % 3) * (lw + 0.06)
    py = TY + 0.26 + (i // 3) * 0.38
    img_placeholder(s, px, py, lw, 0.32, p)

prs.save(PATH)
print("saved; total slides:", len(prs.slides._sldIdLst))
